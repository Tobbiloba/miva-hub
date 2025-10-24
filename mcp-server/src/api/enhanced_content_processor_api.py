#!/usr/bin/env python3
"""
Enhanced MIVA University Content Processing API
FastAPI service with comprehensive error handling, validation, and monitoring
"""

import asyncio
import json
import logging
import queue
import mimetypes
import os
import re
import sys
import tempfile
import time
import uuid
from datetime import datetime, timedelta
from pathlib import Path
from typing import Any, Dict, List, Optional, Union
from uuid import UUID

import psycopg2
from psycopg2.extras import RealDictCursor
from dotenv import load_dotenv
from fastapi import FastAPI, File, UploadFile, Form, HTTPException, BackgroundTasks, Depends, Request, status
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, StreamingResponse
from sse_starlette.sse import EventSourceResponse
from pydantic import BaseModel, Field, field_validator
from slowapi import Limiter, _rate_limit_exceeded_handler
from slowapi.errors import RateLimitExceeded
from slowapi.util import get_remote_address

# Enhanced imports
import sys
import os
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from core.ai_integration import MIVAAIStack
from core.s3_service import s3_service

# Utility functions
def secure_filename(filename: str) -> str:
    """Securely clean a filename to prevent path traversal attacks."""
    if not filename:
        return "upload"
    
    # Remove path components
    filename = os.path.basename(filename)
    
    # Remove dangerous characters, keep only alphanumeric, dots, hyphens, underscores
    filename = re.sub(r'[^a-zA-Z0-9._-]', '_', filename)
    
    # Remove multiple consecutive dots/underscores
    filename = re.sub(r'[._-]{2,}', '_', filename)
    
    # Ensure it doesn't start with a dot
    filename = filename.lstrip('.')
    
    # Limit length
    if len(filename) > 100:
        name, ext = os.path.splitext(filename)
        filename = name[:90] + ext
    
    return filename or "upload"

# Document processing imports
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import whisper
    WHISPER_AVAILABLE = True
except ImportError:
    WHISPER_AVAILABLE = False

# Load environment
load_dotenv()

# Real-time status tracking
import threading
from collections import defaultdict
from typing import Set

class ProcessingStatusTracker:
    """Thread-safe status tracker for real-time updates"""
    def __init__(self):
        self._status_data = {}
        self._subscribers = defaultdict(set)  # job_id -> set of queue objects
        self._lock = threading.Lock()
    
    def update_status(self, job_id: str, status_data: dict):
        """Update status and notify all subscribers"""
        with self._lock:
            self._status_data[job_id] = {
                **status_data,
                'timestamp': datetime.now().isoformat()
            }
            
            # Notify all subscribers
            for queue in self._subscribers[job_id]:
                try:
                    queue.put_nowait(status_data)
                except:
                    pass  # Queue might be full or closed
    
    def subscribe(self, job_id: str, queue):
        """Subscribe to status updates for a job"""
        with self._lock:
            self._subscribers[job_id].add(queue)
    
    def unsubscribe(self, job_id: str, queue):
        """Unsubscribe from status updates"""
        with self._lock:
            self._subscribers[job_id].discard(queue)
    
    def get_current_status(self, job_id: str) -> dict:
        """Get current status for a job"""
        with self._lock:
            return self._status_data.get(job_id, {})

# Global status tracker
status_tracker = ProcessingStatusTracker()

# Processing Queue System
import asyncio
from asyncio import Queue, Lock
from dataclasses import dataclass
from enum import Enum
from typing import Callable, Optional

class JobPriority(Enum):
    LOW = 1
    NORMAL = 2 
    HIGH = 3
    URGENT = 4

@dataclass
class ProcessingJob:
    job_id: str
    material_id: str
    file_path: str
    priority: JobPriority
    created_at: datetime
    retries: int = 0
    max_retries: int = 3
    callback: Optional[Callable] = None

class ProcessingQueue:
    """Async processing queue with priority support and worker management."""
    
    def __init__(self, max_workers: int = 3, max_queue_size: int = 100):
        self.max_workers = max_workers
        self.max_queue_size = max_queue_size
        self.queues = {
            JobPriority.URGENT: Queue(maxsize=max_queue_size),
            JobPriority.HIGH: Queue(maxsize=max_queue_size), 
            JobPriority.NORMAL: Queue(maxsize=max_queue_size),
            JobPriority.LOW: Queue(maxsize=max_queue_size)
        }
        self.workers = []
        self.active_jobs = {}  # job_id -> ProcessingJob
        self.queue_lock = Lock()
        self.stats = {
            'jobs_queued': 0,
            'jobs_completed': 0,
            'jobs_failed': 0,
            'jobs_active': 0,
            'queue_sizes': {}
        }
        self.running = False
    
    async def start(self):
        """Start the queue workers."""
        if self.running:
            return
        
        self.running = True
        for i in range(self.max_workers):
            worker = asyncio.create_task(self._worker(f"worker-{i}"))
            self.workers.append(worker)
        
        logger.info(f"🔧 Processing queue started with {self.max_workers} workers")
    
    async def stop(self):
        """Stop the queue workers gracefully."""
        self.running = False
        
        # Cancel all workers
        for worker in self.workers:
            worker.cancel()
        
        # Wait for workers to finish
        await asyncio.gather(*self.workers, return_exceptions=True)
        self.workers.clear()
        
        logger.info("🛑 Processing queue stopped")
    
    async def add_job(self, job: ProcessingJob) -> bool:
        """Add a job to the appropriate priority queue."""
        try:
            async with self.queue_lock:
                queue = self.queues[job.priority]
                
                # Check if queue is full
                if queue.full():
                    logger.warning(f"Queue for priority {job.priority.name} is full")
                    return False
                
                await queue.put(job)
                self.stats['jobs_queued'] += 1
                
                # Update status tracker
                status_tracker.update_status(job.job_id, {
                    "status": "queued",
                    "stage": "queue",
                    "message": f"Job queued with {job.priority.name} priority",
                    "queue_position": queue.qsize(),
                    "priority": job.priority.name
                })
                
                logger.info(f"📝 Job {job.job_id} queued with {job.priority.name} priority")
                return True
                
        except Exception as e:
            logger.error(f"❌ Failed to queue job {job.job_id}: {str(e)}")
            return False
    
    async def _worker(self, worker_name: str):
        """Worker coroutine that processes jobs from the queues."""
        logger.info(f"👷 Worker {worker_name} started")
        
        while self.running:
            try:
                # Get next job from highest priority queue
                job = await self._get_next_job()
                
                if job is None:
                    await asyncio.sleep(1)  # No jobs available, wait a bit
                    continue
                
                # Process the job
                await self._process_job(job, worker_name)
                
            except asyncio.CancelledError:
                logger.info(f"👷 Worker {worker_name} cancelled")
                break
            except Exception as e:
                logger.error(f"❌ Worker {worker_name} error: {str(e)}")
                await asyncio.sleep(5)  # Wait before retrying
        
        logger.info(f"👷 Worker {worker_name} stopped")
    
    async def _get_next_job(self) -> Optional[ProcessingJob]:
        """Get the next job from the highest priority queue that has jobs."""
        async with self.queue_lock:
            # Check queues in priority order
            for priority in [JobPriority.URGENT, JobPriority.HIGH, JobPriority.NORMAL, JobPriority.LOW]:
                queue = self.queues[priority]
                if not queue.empty():
                    try:
                        job = await asyncio.wait_for(queue.get(), timeout=0.1)
                        self.active_jobs[job.job_id] = job
                        self.stats['jobs_active'] += 1
                        return job
                    except asyncio.TimeoutError:
                        continue
            
            return None
    
    async def _process_job(self, job: ProcessingJob, worker_name: str):
        """Process a single job."""
        try:
            logger.info(f"🔄 {worker_name} processing job {job.job_id}")
            
            # Update status to processing
            status_tracker.update_status(job.job_id, {
                "status": "processing",
                "stage": "dequeued",
                "message": f"Processing started by {worker_name}",
                "worker": worker_name,
                "retries": job.retries
            })
            
            # Process the actual content
            await start_background_processing(job.material_id, job.job_id, job.file_path)
            
            # Job completed successfully
            await self._job_completed(job, worker_name)
            
        except Exception as e:
            logger.error(f"❌ {worker_name} job {job.job_id} failed: {str(e)}")
            await self._job_failed(job, worker_name, str(e))
    
    async def _job_completed(self, job: ProcessingJob, worker_name: str):
        """Handle successful job completion."""
        async with self.queue_lock:
            self.active_jobs.pop(job.job_id, None)
            self.stats['jobs_completed'] += 1
            self.stats['jobs_active'] -= 1
        
        logger.info(f"✅ {worker_name} completed job {job.job_id}")
        
        # Call callback if provided
        if job.callback:
            try:
                await job.callback(job, 'completed')
            except Exception as e:
                logger.error(f"Callback error for job {job.job_id}: {str(e)}")
    
    async def _job_failed(self, job: ProcessingJob, worker_name: str, error: str):
        """Handle job failure with retry logic."""
        async with self.queue_lock:
            self.active_jobs.pop(job.job_id, None)
            self.stats['jobs_active'] -= 1
        
        # Check if we should retry
        if job.retries < job.max_retries:
            job.retries += 1
            logger.warning(f"🔄 Retrying job {job.job_id} (attempt {job.retries}/{job.max_retries})")
            
            # Re-queue with lower priority for retry
            retry_priority = JobPriority.LOW if job.priority != JobPriority.LOW else JobPriority.LOW
            job.priority = retry_priority
            
            await asyncio.sleep(2 ** job.retries)  # Exponential backoff
            await self.add_job(job)
        else:
            # Max retries exceeded
            self.stats['jobs_failed'] += 1
            logger.error(f"❌ Job {job.job_id} failed permanently after {job.retries} retries")
            
            # Call callback if provided
            if job.callback:
                try:
                    await job.callback(job, 'failed', error)
                except Exception as e:
                    logger.error(f"Callback error for failed job {job.job_id}: {str(e)}")
    
    def get_queue_stats(self) -> dict:
        """Get current queue statistics."""
        self.stats['queue_sizes'] = {
            priority.name: queue.qsize() 
            for priority, queue in self.queues.items()
        }
        return self.stats.copy()

# Global processing queue
processing_queue = ProcessingQueue(max_workers=3, max_queue_size=50)

# Duplicate Content Detection System
import hashlib
from difflib import SequenceMatcher
import re

class DuplicateDetector:
    """Advanced duplicate detection with multiple similarity algorithms."""
    
    def __init__(self):
        self.similarity_threshold = 0.85  # 85% similarity threshold
        self.hash_algorithms = ['md5', 'sha256']
    
    def calculate_file_hash(self, content: bytes, algorithm: str = 'sha256') -> str:
        """Calculate hash of file content."""
        hash_func = getattr(hashlib, algorithm)()
        hash_func.update(content)
        return hash_func.hexdigest()
    
    def calculate_content_hash(self, text: str) -> str:
        """Calculate hash of normalized text content."""
        # Normalize text for better comparison
        normalized_text = self.normalize_text(text)
        return hashlib.sha256(normalized_text.encode('utf-8')).hexdigest()
    
    def normalize_text(self, text: str) -> str:
        """Normalize text by removing whitespace, case, and non-essential characters."""
        # Convert to lowercase
        text = text.lower()
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text)
        # Remove punctuation (keeping alphanumeric and spaces)
        text = re.sub(r'[^\w\s]', '', text)
        # Strip leading/trailing whitespace
        return text.strip()
    
    def calculate_text_similarity(self, text1: str, text2: str) -> float:
        """Calculate similarity between two texts using sequence matching."""
        # Normalize both texts
        norm_text1 = self.normalize_text(text1)
        norm_text2 = self.normalize_text(text2)
        
        # Use SequenceMatcher for similarity
        matcher = SequenceMatcher(None, norm_text1, norm_text2)
        return matcher.ratio()
    
    def calculate_semantic_similarity(self, text1: str, text2: str) -> float:
        """Calculate semantic similarity using word overlap and structure."""
        # Normalize texts
        norm_text1 = self.normalize_text(text1)
        norm_text2 = self.normalize_text(text2)
        
        # Split into words
        words1 = set(norm_text1.split())
        words2 = set(norm_text2.split())
        
        # Calculate Jaccard similarity (intersection over union)
        if not words1 and not words2:
            return 1.0
        if not words1 or not words2:
            return 0.0
        
        intersection = words1.intersection(words2)
        union = words1.union(words2)
        
        return len(intersection) / len(union)
    
    async def find_duplicates(self, content: bytes, text_content: str, 
                            filename: str, course_id: str = None) -> dict:
        """Find potential duplicates in the database."""
        duplicates = {
            "exact_file_matches": [],
            "content_matches": [],
            "similar_content": [],
            "is_duplicate": False,
            "similarity_scores": {}
        }
        
        try:
            # Calculate hashes
            file_hash = self.calculate_file_hash(content)
            content_hash = self.calculate_content_hash(text_content)
            
            # Query database for potential matches
            conn = get_db_connection()
            cursor = conn.cursor()
            
            # First check for exact file matches (same hash)
            cursor.execute("""
                SELECT DISTINCT cm.id, cm.title, cm.file_name, cm.file_path, 
                       apc.extracted_text, cm."course_id"
                FROM course_material cm
                LEFT JOIN ai_processed_content apc ON cm.id = apc.course_material_id
                WHERE (%s IS NULL OR cm."course_id" = %s)
                AND cm.file_path IS NOT NULL
            """, (course_id, course_id))
            
            existing_materials = cursor.fetchall()
            cursor.close()
            conn.close()
            
            # Check each existing material
            for material in existing_materials:
                if isinstance(material, dict):
                    material_id = material['id']
                    title = material['title']
                    file_name = material['file_name']
                    file_path = material['file_path']
                    extracted_text = material['extracted_text']
                    material_course_id = material['course_id']
                else:
                    material_id, title, file_name, file_path, extracted_text, material_course_id = material
                
                # Skip if no extracted text
                if not extracted_text:
                    continue
                
                # Check file hash if file exists
                try:
                    if file_path and Path(file_path).exists():
                        with open(file_path, 'rb') as f:
                            existing_content = f.read()
                        existing_file_hash = self.calculate_file_hash(existing_content)
                        
                        if existing_file_hash == file_hash:
                            duplicates["exact_file_matches"].append({
                                "material_id": material_id,
                                "title": title,
                                "filename": file_name,
                                "course_id": material_course_id,
                                "match_type": "exact_file"
                            })
                            duplicates["is_duplicate"] = True
                            continue
                except Exception as e:
                    logger.warning(f"Could not check file hash for {file_path}: {str(e)}")
                
                # Check content hash
                existing_content_hash = self.calculate_content_hash(extracted_text)
                if existing_content_hash == content_hash:
                    duplicates["content_matches"].append({
                        "material_id": material_id,
                        "title": title,
                        "filename": file_name,
                        "course_id": material_course_id,
                        "match_type": "exact_content"
                    })
                    duplicates["is_duplicate"] = True
                    continue
                
                # Calculate similarity scores
                text_similarity = self.calculate_text_similarity(text_content, extracted_text)
                semantic_similarity = self.calculate_semantic_similarity(text_content, extracted_text)
                
                # Average the similarities for overall score
                overall_similarity = (text_similarity + semantic_similarity) / 2
                
                if overall_similarity >= self.similarity_threshold:
                    duplicates["similar_content"].append({
                        "material_id": material_id,
                        "title": title,
                        "filename": file_name,
                        "course_id": material_course_id,
                        "text_similarity": round(text_similarity, 3),
                        "semantic_similarity": round(semantic_similarity, 3),
                        "overall_similarity": round(overall_similarity, 3),
                        "match_type": "similar_content"
                    })
                    
                    duplicates["similarity_scores"][material_id] = {
                        "text_similarity": text_similarity,
                        "semantic_similarity": semantic_similarity,
                        "overall_similarity": overall_similarity
                    }
                    
                    # Mark as duplicate if very high similarity
                    if overall_similarity >= 0.95:
                        duplicates["is_duplicate"] = True
            
            return duplicates
            
        except Exception as e:
            logger.error(f"❌ Duplicate detection failed: {str(e)}")
            return {
                "exact_file_matches": [],
                "content_matches": [],
                "similar_content": [],
                "is_duplicate": False,
                "similarity_scores": {},
                "error": str(e)
            }

# Global duplicate detector
duplicate_detector = DuplicateDetector()

# Enhanced logging configuration
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('content_processor.log'),
        logging.StreamHandler(sys.stdout)
    ]
)
logger = logging.getLogger(__name__)

# Rate limiting
limiter = Limiter(key_func=get_remote_address)

# Configuration
class Config:
    MAX_FILE_SIZE = int(os.getenv('MAX_FILE_SIZE', 100 * 1024 * 1024))  # 100MB
    MAX_FILES_PER_HOUR = int(os.getenv('MAX_FILES_PER_HOUR', 50))
    ALLOWED_EXTENSIONS = {'.pdf', '.txt', '.docx', '.pptx', '.mp4', '.avi', '.mov', '.mp3', '.wav'}
    UPLOAD_DIR = Path(os.getenv('UPLOAD_DIR', 'uploads'))
    UPLOAD_DIR.mkdir(exist_ok=True)

# Database configuration - unified with Next.js frontend
def get_db_config():
    """Get unified database configuration using POSTGRES_URL"""
    postgres_url = os.getenv('POSTGRES_URL') or os.getenv('DATABASE_URL')
    
    if postgres_url:
        # Parse POSTGRES_URL
        import urllib.parse as urlparse
        parsed = urlparse.urlparse(postgres_url)
        return {
            'host': parsed.hostname,
            'port': parsed.port or 5432,
            'database': parsed.path.lstrip('/'),
            'user': parsed.username,
            'password': parsed.password,
            'cursor_factory': RealDictCursor
        }
    else:
        # Fallback to individual environment variables
        return {
            'host': os.getenv('DB_HOST', 'localhost'),
            'port': int(os.getenv('DB_PORT', '5432')),
            'database': os.getenv('DB_NAME', 'miva_hub'),  # Changed from miva_academic
            'user': os.getenv('DB_USER', 'postgres'),
            'password': os.getenv('DB_PASSWORD', ''),
            'cursor_factory': RealDictCursor
        }

DB_CONFIG = get_db_config()

# Custom Exception Classes
class ContentProcessingError(Exception):
    """Base exception for content processing errors"""
    def __init__(self, message: str, error_code: str = "PROCESSING_ERROR", details: Dict = None):
        self.message = message
        self.error_code = error_code
        self.details = details or {}
        super().__init__(self.message)

class FileValidationError(ContentProcessingError):
    """Exception for file validation errors"""
    def __init__(self, message: str, details: Dict = None):
        super().__init__(message, "FILE_VALIDATION_ERROR", details)

class DatabaseError(ContentProcessingError):
    """Exception for database-related errors"""
    def __init__(self, message: str, details: Dict = None):
        super().__init__(message, "DATABASE_ERROR", details)

class AIProcessingError(ContentProcessingError):
    """Exception for AI processing errors"""
    def __init__(self, message: str, details: Dict = None):
        super().__init__(message, "AI_PROCESSING_ERROR", details)

# Database Helper Functions
def get_db_connection():
    """Get database connection with enhanced error handling"""
    try:
        return psycopg2.connect(**DB_CONFIG)
    except Exception as e:
        logger.error(f"❌ Database connection failed: {e}")
        raise DatabaseError("Database service temporarily unavailable")

# Enhanced Content Processors

class EnhancedPDFProcessor:
    """Advanced PDF processor with metadata extraction and structure analysis"""
    
    def __init__(self):
        self.logger = logging.getLogger(f"{__name__}.PDFProcessor")
    
    async def process_file(self, file_path: str) -> str:
        """Process PDF file with enhanced metadata and structure extraction"""
        try:
            self.logger.info(f"📄 Processing PDF: {file_path}")
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                # Extract metadata
                metadata = self._extract_metadata(pdf_reader)
                
                # Extract text with page information
                pages_content = []
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text.strip():
                        pages_content.append({
                            'page': page_num + 1,
                            'content': page_text.strip()
                        })
                
                # Combine content with structure
                content_parts = [
                    "=== PDF DOCUMENT ANALYSIS ===",
                    f"Title: {metadata.get('title', 'Unknown')}",
                    f"Author: {metadata.get('author', 'Unknown')}",
                    f"Pages: {len(pdf_reader.pages)}",
                    f"Creation Date: {metadata.get('creation_date', 'Unknown')}",
                    "",
                    "=== DOCUMENT CONTENT ===",
                ]
                
                # Add page-by-page content
                for page_info in pages_content:
                    content_parts.append(f"\n--- Page {page_info['page']} ---")
                    content_parts.append(page_info['content'])
                
                full_content = '\n'.join(content_parts)
                
                self.logger.info(f"✅ PDF processed: {len(pages_content)} pages, {len(full_content)} characters")
                return full_content
                
        except Exception as e:
            self.logger.error(f"❌ PDF processing failed: {str(e)}")
            raise ContentProcessingError(f"PDF processing failed: {str(e)}")
    
    def _extract_metadata(self, pdf_reader) -> Dict[str, Any]:
        """Extract PDF metadata"""
        try:
            metadata = {}
            if pdf_reader.metadata:
                metadata['title'] = pdf_reader.metadata.get('/Title', '')
                metadata['author'] = pdf_reader.metadata.get('/Author', '')
                metadata['subject'] = pdf_reader.metadata.get('/Subject', '')
                metadata['creator'] = pdf_reader.metadata.get('/Creator', '')
                metadata['producer'] = pdf_reader.metadata.get('/Producer', '')
                
                # Handle creation date
                creation_date = pdf_reader.metadata.get('/CreationDate', '')
                if creation_date:
                    metadata['creation_date'] = str(creation_date)
                    
            return metadata
        except Exception as e:
            self.logger.warning(f"⚠️ Metadata extraction failed: {str(e)}")
            return {}

class EnhancedDOCXProcessor:
    """Advanced DOCX processor with styles, formatting, and structure analysis"""
    
    def __init__(self):
        self.logger = logging.getLogger(f"{__name__}.DOCXProcessor")
    
    async def process_file(self, file_path: str) -> str:
        """Process DOCX file with enhanced structure and formatting extraction"""
        try:
            self.logger.info(f"📝 Processing DOCX: {file_path}")
            
            doc = docx.Document(file_path)
            
            # Extract document properties
            properties = self._extract_properties(doc)
            
            # Extract content with structure
            content_parts = [
                "=== DOCX DOCUMENT ANALYSIS ===",
                f"Title: {properties.get('title', 'Unknown')}",
                f"Author: {properties.get('author', 'Unknown')}",
                f"Paragraphs: {len(doc.paragraphs)}",
                f"Created: {properties.get('created', 'Unknown')}",
                "",
                "=== DOCUMENT STRUCTURE ===",
            ]
            
            # Process paragraphs with style information
            for i, paragraph in enumerate(doc.paragraphs):
                if paragraph.text.strip():
                    style_name = paragraph.style.name if paragraph.style else "Normal"
                    
                    # Identify headings and important content
                    if any(heading in style_name.lower() for heading in ['heading', 'title']):
                        content_parts.append(f"\n### {paragraph.text.strip()} ###")
                    elif style_name == "Normal":
                        content_parts.append(paragraph.text.strip())
                    else:
                        content_parts.append(f"[{style_name}] {paragraph.text.strip()}")
            
            # Extract tables if any
            if doc.tables:
                content_parts.append("\n=== DOCUMENT TABLES ===")
                for table_idx, table in enumerate(doc.tables):
                    content_parts.append(f"\n--- Table {table_idx + 1} ---")
                    table_content = self._extract_table_content(table)
                    content_parts.append(table_content)
            
            full_content = '\n'.join(content_parts)
            
            self.logger.info(f"✅ DOCX processed: {len(doc.paragraphs)} paragraphs, {len(doc.tables)} tables")
            return full_content
            
        except Exception as e:
            self.logger.error(f"❌ DOCX processing failed: {str(e)}")
            raise ContentProcessingError(f"DOCX processing failed: {str(e)}")
    
    def _extract_properties(self, doc) -> Dict[str, Any]:
        """Extract document properties"""
        try:
            properties = {}
            if hasattr(doc, 'core_properties'):
                cp = doc.core_properties
                properties['title'] = cp.title or ''
                properties['author'] = cp.author or ''
                properties['subject'] = cp.subject or ''
                properties['created'] = str(cp.created) if cp.created else ''
                properties['modified'] = str(cp.modified) if cp.modified else ''
            return properties
        except Exception as e:
            self.logger.warning(f"⚠️ Properties extraction failed: {str(e)}")
            return {}
    
    def _extract_table_content(self, table) -> str:
        """Extract table content in readable format"""
        try:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):  # Only add non-empty rows
                    rows.append(" | ".join(cells))
            return '\n'.join(rows)
        except Exception as e:
            return f"Table extraction error: {str(e)}"

class EnhancedPPTXProcessor:
    """Advanced PPTX processor with slide-level content and layout analysis"""
    
    def __init__(self):
        self.logger = logging.getLogger(f"{__name__}.PPTXProcessor")
    
    async def process_file(self, file_path: str) -> str:
        """Process PPTX file with slide-by-slide content extraction"""
        try:
            self.logger.info(f"🎯 Processing PPTX: {file_path}")
            
            prs = Presentation(file_path)
            
            # Extract presentation properties
            properties = self._extract_properties(prs)
            
            content_parts = [
                "=== POWERPOINT PRESENTATION ANALYSIS ===",
                f"Title: {properties.get('title', 'Unknown')}",
                f"Author: {properties.get('author', 'Unknown')}",
                f"Slides: {len(prs.slides)}",
                f"Created: {properties.get('created', 'Unknown')}",
                "",
                "=== SLIDE CONTENT ===",
            ]
            
            # Process each slide
            for slide_idx, slide in enumerate(prs.slides):
                content_parts.append(f"\n--- Slide {slide_idx + 1} ---")
                
                # Extract slide layout information
                layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown Layout"
                content_parts.append(f"Layout: {layout_name}")
                
                # Extract text from all shapes
                slide_text = self._extract_slide_text(slide)
                if slide_text:
                    content_parts.append(slide_text)
                else:
                    content_parts.append("[No text content on this slide]")
                
                # Extract notes if any
                if hasattr(slide, 'notes_slide') and slide.notes_slide:
                    notes_text = self._extract_notes_text(slide.notes_slide)
                    if notes_text:
                        content_parts.append(f"Speaker Notes: {notes_text}")
            
            full_content = '\n'.join(content_parts)
            
            self.logger.info(f"✅ PPTX processed: {len(prs.slides)} slides")
            return full_content
            
        except Exception as e:
            self.logger.error(f"❌ PPTX processing failed: {str(e)}")
            raise ContentProcessingError(f"PPTX processing failed: {str(e)}")
    
    def _extract_properties(self, prs) -> Dict[str, Any]:
        """Extract presentation properties"""
        try:
            properties = {}
            if hasattr(prs, 'core_properties'):
                cp = prs.core_properties
                properties['title'] = cp.title or ''
                properties['author'] = cp.author or ''
                properties['subject'] = cp.subject or ''
                properties['created'] = str(cp.created) if cp.created else ''
                properties['modified'] = str(cp.modified) if cp.modified else ''
            return properties
        except Exception as e:
            self.logger.warning(f"⚠️ Properties extraction failed: {str(e)}")
            return {}
    
    def _extract_slide_text(self, slide) -> str:
        """Extract all text from slide shapes"""
        text_parts = []
        try:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    # Identify shape type for better formatting
                    if hasattr(shape, 'name') and 'title' in shape.name.lower():
                        text_parts.append(f"TITLE: {shape.text.strip()}")
                    else:
                        text_parts.append(shape.text.strip())
                        
                # Extract table content if shape is a table
                if hasattr(shape, 'table') and shape.table:
                    table_content = self._extract_table_content(shape.table)
                    if table_content:
                        text_parts.append(f"TABLE:\n{table_content}")
                        
        except Exception as e:
            self.logger.warning(f"⚠️ Slide text extraction failed: {str(e)}")
            
        return '\n'.join(text_parts)
    
    def _extract_notes_text(self, notes_slide) -> str:
        """Extract speaker notes text"""
        try:
            notes_text = []
            for shape in notes_slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    notes_text.append(shape.text.strip())
            return ' '.join(notes_text)
        except Exception as e:
            return ""
    
    def _extract_table_content(self, table) -> str:
        """Extract table content from PPTX table"""
        try:
            rows = []
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    cells.append(cell.text.strip())
                if any(cells):
                    rows.append(" | ".join(cells))
            return '\n'.join(rows)
        except Exception as e:
            return f"Table extraction error: {str(e)}"

class EnhancedAudioVideoProcessor:
    """Advanced audio/video processor with Whisper AI transcription"""
    
    def __init__(self):
        self.logger = logging.getLogger(f"{__name__}.AudioVideoProcessor")
        self.whisper_model = None
    
    async def process_file(self, file_path: str) -> str:
        """Process audio/video file with AI transcription"""
        try:
            self.logger.info(f"🎵 Processing audio/video: {file_path}")
            
            if not WHISPER_AVAILABLE:
                raise ContentProcessingError("Whisper not available for audio/video processing")
            
            # Load Whisper model if not already loaded
            if self.whisper_model is None:
                self.logger.info("🔄 Loading Whisper model...")
                import whisper
                self.whisper_model = whisper.load_model("base")  # Use base model for speed
                
            # Transcribe audio
            self.logger.info("🎯 Starting transcription...")
            result = self.whisper_model.transcribe(file_path)
            
            # Extract detailed information
            content_parts = [
                "=== AUDIO/VIDEO TRANSCRIPTION ===",
                f"File: {Path(file_path).name}",
                f"Language: {result.get('language', 'Unknown')}",
                "",
                "=== TRANSCRIPT ===",
                result['text']
            ]
            
            # Add segments if available for timestamped content
            if 'segments' in result and result['segments']:
                content_parts.append("\n=== TIMESTAMPED SEGMENTS ===")
                for segment in result['segments'][:10]:  # Limit to first 10 segments
                    start_time = self._format_timestamp(segment.get('start', 0))
                    end_time = self._format_timestamp(segment.get('end', 0))
                    text = segment.get('text', '').strip()
                    content_parts.append(f"[{start_time} - {end_time}] {text}")
            
            full_content = '\n'.join(content_parts)
            
            self.logger.info(f"✅ Audio/video processed: {len(result['text'])} characters transcribed")
            return full_content
            
        except Exception as e:
            self.logger.error(f"❌ Audio/video processing failed: {str(e)}")
            raise ContentProcessingError(f"Audio/video processing failed: {str(e)}")
    
    def _format_timestamp(self, seconds: float) -> str:
        """Format timestamp in MM:SS format"""
        minutes = int(seconds // 60)
        seconds = int(seconds % 60)
        return f"{minutes:02d}:{seconds:02d}"

# Initialize enhanced processors
enhanced_pdf_processor = EnhancedPDFProcessor()
enhanced_docx_processor = EnhancedDOCXProcessor()
enhanced_pptx_processor = EnhancedPPTXProcessor()
enhanced_audio_video_processor = EnhancedAudioVideoProcessor()

# Pydantic Models for Request/Response Validation
class ContentUploadRequest(BaseModel):
    title: str = Field(..., min_length=1, max_length=200, description="Content title")
    course_id: Optional[str] = Field(None, description="Course ID (UUID)")
    week_number: Optional[int] = Field(None, ge=1, le=20, description="Week number")
    material_type: Optional[str] = Field(None, pattern="^(pdf|video|audio|text|interactive)$")
    
    @field_validator('title')
    @classmethod
    def validate_title(cls, v):
        if not v.strip():
            raise ValueError('Title cannot be empty')
        return v.strip()

class ProcessingStatusResponse(BaseModel):
    processing_id: str
    material_id: str  # UUID
    job_id: str  # UUID
    status: str = Field(..., pattern="^(pending|processing|completed|failed)$")
    progress: int = Field(..., ge=0, le=100)
    error_message: Optional[str] = None
    ai_processed: bool
    embedding_count: int
    created_at: Optional[datetime] = None
    started_at: Optional[datetime] = None
    completed_at: Optional[datetime] = None
    estimated_time_remaining: Optional[int] = None  # seconds

class SearchRequest(BaseModel):
    query: str = Field(..., min_length=1, max_length=500)
    course_filter: Optional[str] = Field(None, max_length=10)
    similarity_threshold: float = Field(0.3, ge=0.0, le=1.0)
    limit: int = Field(10, ge=1, le=50)

class SearchResult(BaseModel):
    course_material_id: str  # UUID
    material_title: str
    material_type: str
    chunk_text: str
    similarity_score: float
    course_code: str

class SearchResponse(BaseModel):
    query: str
    results: List[SearchResult]
    count: int
    execution_time_ms: int

class HealthCheckResponse(BaseModel):
    status: str
    timestamp: datetime
    version: str
    services: Dict[str, str]
    system_info: Dict[str, Any]
    performance_metrics: Dict[str, float]

class SystemMetrics(BaseModel):
    total_uploads: int
    successful_uploads: int
    failed_uploads: int
    total_processing_time: float
    average_processing_time: float
    active_jobs: int
    queue_size: int

# Initialize FastAPI app with enhanced configuration
app = FastAPI(
    title="Enhanced MIVA University Content Processor",
    description="Production-ready AI-powered content processing service",
    version="2.0.0",
    docs_url="/docs",
    redoc_url="/redoc"
)

# Add rate limiting
app.state.limiter = limiter
app.add_exception_handler(RateLimitExceeded, _rate_limit_exceeded_handler)

# CORS configuration
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://localhost:3001", "*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Global AI stack instance
ai_stack = None

# Supported MIME types with enhanced detection
SUPPORTED_MIME_TYPES = {
    'application/pdf': 'pdf',
    'text/plain': 'text',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
    'video/mp4': 'video',
    'video/avi': 'video',
    'video/quicktime': 'video',
    'audio/mp3': 'audio',
    'audio/wav': 'audio',
    'audio/m4a': 'audio',
    'image/jpeg': 'image',
    'image/png': 'image',
}

# Processing status tracking
class ProcessingStatus:
    PENDING = "pending"
    PROCESSING = "processing"
    COMPLETED = "completed"
    FAILED = "failed"

# Utility Functions
def get_db_connection():
    """Get database connection with enhanced error handling"""
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        return conn
    except psycopg2.Error as e:
        logger.error(f"Database connection failed: {e}")
        raise DatabaseError(f"Database connection failed: {str(e)}")

async def get_ai_stack():
    """Get AI stack instance with connection validation"""
    global ai_stack
    if ai_stack is None:
        try:
            logger.info("🔄 Initializing AI stack...")
            ai_stack = MIVAAIStack()
            logger.info("🔗 Testing AI stack connection...")
            connection_test = await ai_stack.test_connection()
            logger.info(f"🔗 AI stack connection test result: {connection_test}")
            if not connection_test:
                logger.error("❌ AI stack connection test failed")
                raise AIProcessingError("AI services unavailable - connection test failed")
            logger.info("✅ AI stack initialized and connected successfully")
        except Exception as e:
            logger.error(f"❌ AI stack initialization failed: {e}")
            logger.error(f"❌ Exception type: {type(e).__name__}")
            logger.error(f"❌ Exception details: {str(e)}")
            raise AIProcessingError(f"AI services unavailable: {str(e)}")
    else:
        logger.info("✅ Using existing AI stack instance")
    return ai_stack

def validate_file(file: UploadFile) -> tuple[str, str]:
    """Enhanced file validation with detailed error reporting"""
    try:
        # Check file size
        if file.size and file.size > Config.MAX_FILE_SIZE:
            raise FileValidationError(
                f"File size {file.size} exceeds maximum allowed size {Config.MAX_FILE_SIZE}",
                {"file_size": file.size, "max_size": Config.MAX_FILE_SIZE}
            )
        
        # Check file extension
        file_extension = Path(file.filename).suffix.lower()
        if file_extension not in Config.ALLOWED_EXTENSIONS:
            raise FileValidationError(
                f"File extension '{file_extension}' not supported",
                {"file_extension": file_extension, "allowed_extensions": list(Config.ALLOWED_EXTENSIONS)}
            )
        
        # Determine content type
        content_type = file.content_type or mimetypes.guess_type(file.filename)[0]
        if content_type not in SUPPORTED_MIME_TYPES:
            raise FileValidationError(
                f"Content type '{content_type}' not supported",
                {"content_type": content_type, "supported_types": list(SUPPORTED_MIME_TYPES.keys())}
            )
        
        file_type = SUPPORTED_MIME_TYPES[content_type]
        logger.info(f"File validation successful: {file.filename} ({file_type})")
        return file_type, content_type
        
    except FileValidationError:
        raise
    except Exception as e:
        logger.error(f"File validation error: {e}")
        raise FileValidationError(f"File validation failed: {str(e)}")

# Enhanced Content Processors
class EnhancedContentProcessor:
    """Enhanced content processor with improved extraction"""
    
    @staticmethod
    def create_content_chunks(text: str, chunk_size: int = 300, overlap: int = 50) -> List[str]:
        """Create overlapping text chunks with better context preservation"""
        if not text.strip():
            return []
        
        # Split by sentences first for better chunk boundaries
        sentences = text.replace('\n', ' ').split('. ')
        words = []
        for sentence in sentences:
            words.extend(sentence.split())
        
        chunks = []
        for i in range(0, len(words), chunk_size - overlap):
            chunk = ' '.join(words[i:i + chunk_size])
            if len(chunk.strip()) > 50:  # Only meaningful chunks
                chunks.append(chunk)
        
        return chunks

class EnhancedPDFProcessor(EnhancedContentProcessor):
    """Enhanced PDF processor with better text extraction"""
    
    @staticmethod
    async def extract_text(file_path: str) -> str:
        """Extract text from PDF with enhanced error handling"""
        if not PDF_AVAILABLE:
            raise AIProcessingError("PDF processing not available")
        
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            text += f"Page {page_num + 1}:\n{page_text}\n\n"
                    except Exception as e:
                        logger.warning(f"Failed to extract text from page {page_num + 1}: {e}")
                        continue
            
            if not text.strip():
                raise AIProcessingError("No text could be extracted from PDF")
            
            return text.strip()
        except Exception as e:
            logger.error(f"PDF text extraction failed: {e}")
            raise AIProcessingError(f"PDF text extraction failed: {str(e)}")

class EnhancedDOCXProcessor(EnhancedContentProcessor):
    """Enhanced DOCX processor"""
    
    @staticmethod
    async def extract_text(file_path: str) -> str:
        """Extract text from DOCX file"""
        if not DOCX_AVAILABLE:
            raise AIProcessingError("DOCX processing not available")
        
        try:
            doc = docx.Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            
            if not text.strip():
                raise AIProcessingError("No text could be extracted from DOCX")
            
            return text.strip()
        except Exception as e:
            logger.error(f"DOCX text extraction failed: {e}")
            raise AIProcessingError(f"DOCX text extraction failed: {str(e)}")

class EnhancedPPTXProcessor(EnhancedContentProcessor):
    """Enhanced PPTX processor"""
    
    @staticmethod
    async def extract_text(file_path: str) -> str:
        """Extract text from PPTX file"""
        if not PPTX_AVAILABLE:
            raise AIProcessingError("PPTX processing not available")
        
        try:
            prs = Presentation(file_path)
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"Slide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                text += slide_text + "\n"
            
            if not text.strip():
                raise AIProcessingError("No text could be extracted from PPTX")
            
            return text.strip()
        except Exception as e:
            logger.error(f"PPTX text extraction failed: {e}")
            raise AIProcessingError(f"PPTX text extraction failed: {str(e)}")

class EnhancedTextProcessor(EnhancedContentProcessor):
    """Enhanced text processor"""
    
    @staticmethod
    async def extract_text(file_path: str) -> str:
        """Extract text from plain text file with encoding detection"""
        try:
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        text = file.read()
                        if text.strip():
                            return text.strip()
                except UnicodeDecodeError:
                    continue
            
            raise AIProcessingError("Could not decode text file with any supported encoding")
        except Exception as e:
            logger.error(f"Text extraction failed: {e}")
            raise AIProcessingError(f"Text extraction failed: {str(e)}")

# Exception Handlers
@app.exception_handler(ContentProcessingError)
async def content_processing_exception_handler(request: Request, exc: ContentProcessingError):
    """Handle content processing exceptions"""
    logger.error(f"Content processing error: {exc.message}", extra={"error_code": exc.error_code, "details": exc.details})
    return JSONResponse(
        status_code=400,
        content={
            "error": exc.error_code,
            "message": exc.message,
            "details": exc.details,
            "timestamp": datetime.now().isoformat()
        }
    )

@app.exception_handler(FileValidationError)
async def file_validation_exception_handler(request: Request, exc: FileValidationError):
    """Handle file validation exceptions"""
    logger.warning(f"File validation error: {exc.message}", extra={"details": exc.details})
    return JSONResponse(
        status_code=422,
        content={
            "error": exc.error_code,
            "message": exc.message,
            "details": exc.details,
            "timestamp": datetime.now().isoformat()
        }
    )

@app.exception_handler(DatabaseError)
async def database_exception_handler(request: Request, exc: DatabaseError):
    """Handle database exceptions"""
    logger.error(f"Database error: {exc.message}", extra={"details": exc.details})
    return JSONResponse(
        status_code=503,
        content={
            "error": exc.error_code,
            "message": "Database service temporarily unavailable",
            "timestamp": datetime.now().isoformat()
        }
    )

@app.exception_handler(AIProcessingError)
async def ai_processing_exception_handler(request: Request, exc: AIProcessingError):
    """Handle AI processing exceptions"""
    logger.error(f"AI processing error: {exc.message}", extra={"details": exc.details})
    return JSONResponse(
        status_code=503,
        content={
            "error": exc.error_code,
            "message": "AI processing service temporarily unavailable",
            "timestamp": datetime.now().isoformat()
        }
    )

# Startup Event
@app.on_event("startup")
async def startup_event():
    """Initialize services on startup with comprehensive checks"""
    logger.info("🚀 Enhanced MIVA Content Processor starting up...")
    
    # Test database connection
    try:
        conn = get_db_connection()
        conn.close()
        logger.info("✅ Database connection verified")
    except Exception as e:
        logger.error(f"❌ Database connection failed: {e}")
        raise
    
    # Test AI stack
    try:
        stack = await get_ai_stack()
        logger.info("✅ AI stack initialized and verified")
    except Exception as e:
        logger.error(f"❌ AI stack initialization failed: {e}")
        raise
    
    # Start processing queue
    try:
        await processing_queue.start()
        logger.info("✅ Processing queue started")
    except Exception as e:
        logger.error(f"❌ Processing queue startup failed: {e}")
        raise
    
    # Create upload directory
    Config.UPLOAD_DIR.mkdir(exist_ok=True)
    logger.info(f"✅ Upload directory ready: {Config.UPLOAD_DIR}")
    
    logger.info("🎉 Enhanced MIVA Content Processor ready for requests!")


@app.on_event("shutdown")
async def shutdown_event():
    """Gracefully shutdown services."""
    logger.info("🛑 Shutting down Enhanced MIVA Content Processor...")
    
    try:
        await processing_queue.stop()
        logger.info("✅ Processing queue stopped gracefully")
    except Exception as e:
        logger.error(f"❌ Error stopping processing queue: {e}")
    
    logger.info("👋 Enhanced MIVA Content Processor shutdown complete")

# Enhanced API Endpoints
@app.get("/", response_model=Dict[str, Any])
async def root():
    """Enhanced root endpoint with detailed service information"""
    return {
        "service": "Enhanced MIVA University Content Processor",
        "version": "2.0.0",
        "status": "operational",
        "features": {
            "ai_processing": True,
            "semantic_search": True,
            "multi_format_support": True,
            "rate_limiting": True,
            "enhanced_monitoring": True
        },
        "supported_formats": list(SUPPORTED_MIME_TYPES.keys()),
        "processors": {
            "pdf": PDF_AVAILABLE,
            "docx": DOCX_AVAILABLE,
            "pptx": PPTX_AVAILABLE,
            "audio_video": WHISPER_AVAILABLE
        },
        "timestamp": datetime.now().isoformat()
    }

@app.get("/health", response_model=HealthCheckResponse)
async def health_check():
    """Comprehensive health check with system metrics"""
    start_time = time.time()
    
    # Test database
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute("SELECT 1")
        cursor.fetchone()
        cursor.close()
        conn.close()
        db_status = "healthy"
        db_response_time = time.time() - start_time
    except Exception as e:
        db_status = f"unhealthy: {str(e)}"
        db_response_time = -1
    
    # Test AI stack
    try:
        stack = await get_ai_stack()
        ai_status = "healthy"
        ai_response_time = time.time() - start_time - db_response_time
    except Exception as e:
        ai_status = f"unhealthy: {str(e)}"
        ai_response_time = -1
    
    # System information
    system_info = {
        "upload_directory": str(Config.UPLOAD_DIR),
        "max_file_size_mb": Config.MAX_FILE_SIZE // (1024 * 1024),
        "supported_extensions": list(Config.ALLOWED_EXTENSIONS),
        "python_version": sys.version.split()[0],
        "platform": sys.platform
    }
    
    # Performance metrics
    performance_metrics = {
        "db_response_time_ms": round(db_response_time * 1000, 2) if db_response_time > 0 else -1,
        "ai_response_time_ms": round(ai_response_time * 1000, 2) if ai_response_time > 0 else -1,
        "total_response_time_ms": round((time.time() - start_time) * 1000, 2)
    }
    
    overall_status = "healthy" if db_status == "healthy" and ai_status == "healthy" else "degraded"
    
    return HealthCheckResponse(
        status=overall_status,
        timestamp=datetime.now(),
        version="2.0.0",
        services={
            "database": db_status,
            "ai_stack": ai_status,
            "file_storage": "healthy" if Config.UPLOAD_DIR.exists() else "unhealthy"
        },
        system_info=system_info,
        performance_metrics=performance_metrics
    )

# Enhanced file validation functions
async def validate_file_enhanced(file: UploadFile) -> None:
    """Enhanced file validation with detailed error reporting."""
    if not file.filename:
        raise FileValidationError("No filename provided")
    
    # Check file size
    if hasattr(file, 'size') and file.size and file.size > Config.MAX_FILE_SIZE:
        max_size_mb = Config.MAX_FILE_SIZE / 1024 / 1024
        raise FileValidationError(
            f"File size ({file.size / 1024 / 1024:.1f}MB) exceeds maximum allowed size ({max_size_mb:.0f}MB)",
            details={"max_size_mb": max_size_mb, "file_size_mb": file.size / 1024 / 1024}
        )
    
    # Check file extension
    file_extension = Path(file.filename).suffix.lower()
    if file_extension not in Config.ALLOWED_EXTENSIONS:
        raise FileValidationError(
            f"File type '{file_extension}' not supported",
            details={"supported_extensions": list(Config.ALLOWED_EXTENSIONS)}
        )
    
    # Check MIME type
    content_type = file.content_type
    if content_type and not any(mime in content_type for mime in [
        'text/', 'application/pdf', 'application/vnd.openxmlformats-officedocument',
        'video/', 'audio/', 'application/msword'
    ]):
        raise FileValidationError(
            f"MIME type '{content_type}' not supported",
            details={"content_type": content_type}
        )
    
    # Read a small portion to verify it's not corrupted
    try:
        await file.seek(0)
        sample = await file.read(1024)
        await file.seek(0)
        
        if not sample:
            raise FileValidationError("File appears to be empty")
            
    except Exception as e:
        raise FileValidationError(f"Failed to read file: {str(e)}")


async def validate_file_security(file: UploadFile) -> dict:
    """Advanced security validation with malware detection and content analysis."""
    validation_results = {
        "is_safe": True,
        "security_score": 100,
        "warnings": [],
        "threats_detected": [],
        "file_metadata": {}
    }
    
    try:
        # Basic file validation first
        await validate_file_enhanced(file)
        
        # Read file content for analysis
        await file.seek(0)
        content = await file.read()
        await file.seek(0)
        
        # File metadata extraction
        validation_results["file_metadata"] = {
            "size_bytes": len(content),
            "extension": Path(file.filename).suffix.lower(),
            "content_type": file.content_type,
            "filename": secure_filename(file.filename)
        }
        
        # 1. Magic byte validation (file signature check)
        if not await validate_file_signature(content, file.filename):
            validation_results["warnings"].append("File signature doesn't match extension")
            validation_results["security_score"] -= 20
        
        # 2. Suspicious pattern detection
        suspicious_patterns = await detect_suspicious_patterns(content)
        if suspicious_patterns:
            validation_results["threats_detected"].extend(suspicious_patterns)
            validation_results["security_score"] -= len(suspicious_patterns) * 15
        
        # 3. Embedded content analysis for documents
        if file.filename.lower().endswith(('.pdf', '.docx', '.pptx')):
            embedded_threats = await analyze_embedded_content(content, file.filename)
            if embedded_threats:
                validation_results["threats_detected"].extend(embedded_threats)
                validation_results["security_score"] -= len(embedded_threats) * 25
        
        # 4. File size anomaly detection
        expected_size_range = get_expected_file_size_range(file.filename)
        if len(content) < expected_size_range[0] or len(content) > expected_size_range[1]:
            validation_results["warnings"].append(f"Unusual file size for {Path(file.filename).suffix} file")
            validation_results["security_score"] -= 10
        
        # 5. Content entropy analysis (detect compressed/encrypted content)
        entropy = calculate_file_entropy(content)
        if entropy > 7.5:  # High entropy might indicate compression or encryption
            validation_results["warnings"].append(f"High entropy content detected ({entropy:.2f})")
            validation_results["security_score"] -= 15
        
        # 6. Filename security check
        if not is_filename_safe(file.filename):
            validation_results["warnings"].append("Potentially unsafe filename patterns")
            validation_results["security_score"] -= 10
        
        # Determine overall safety
        validation_results["is_safe"] = (
            validation_results["security_score"] >= 60 and
            len(validation_results["threats_detected"]) == 0
        )
        
        # Log security analysis
        if validation_results["security_score"] < 80:
            logger.warning(f"Security concerns for file {file.filename}: Score {validation_results['security_score']}")
        
        return validation_results
        
    except Exception as e:
        logger.error(f"Security validation failed for {file.filename}: {str(e)}")
        return {
            "is_safe": False,
            "security_score": 0,
            "warnings": ["Security validation failed"],
            "threats_detected": [f"Validation error: {str(e)}"],
            "file_metadata": {}
        }


async def validate_file_signature(content: bytes, filename: str) -> bool:
    """Validate file magic bytes match the expected format."""
    if len(content) < 16:
        return False
    
    file_ext = Path(filename).suffix.lower()
    magic_bytes = content[:16]
    
    # Common file signatures
    signatures = {
        '.pdf': [b'%PDF'],
        '.docx': [b'PK\x03\x04'],  # ZIP-based format
        '.pptx': [b'PK\x03\x04'],  # ZIP-based format
        '.zip': [b'PK\x03\x04', b'PK\x05\x06', b'PK\x07\x08'],
        '.mp4': [b'\x00\x00\x00\x18ftypmp4', b'\x00\x00\x00\x1cftypmp42'],
        '.mp3': [b'ID3', b'\xff\xfb', b'\xff\xf3', b'\xff\xf2'],
        '.wav': [b'RIFF'],
        '.txt': [],  # Text files don't have consistent magic bytes
    }
    
    if file_ext not in signatures:
        return True  # Unknown extension, can't validate
    
    expected_signatures = signatures[file_ext]
    if not expected_signatures:  # No specific signature required
        return True
    
    return any(magic_bytes.startswith(sig) for sig in expected_signatures)


async def detect_suspicious_patterns(content: bytes) -> list:
    """Detect suspicious patterns that might indicate malware or malicious content."""
    threats = []
    
    # Convert to string for pattern matching (ignore decode errors)
    try:
        text_content = content.decode('utf-8', errors='ignore').lower()
    except:
        text_content = str(content).lower()
    
    # Suspicious patterns to detect
    suspicious_patterns = {
        'javascript_injection': [
            '<script', 'javascript:', 'eval(', 'document.write',
            'window.location', 'document.cookie'
        ],
        'sql_injection': [
            'union select', 'drop table', 'exec(', '--',
            'xp_cmdshell', 'sp_executesql'
        ],
        'shell_commands': [
            '/bin/sh', '/bin/bash', 'cmd.exe', 'powershell',
            'system(', 'exec(', 'passthru('
        ],
        'suspicious_urls': [
            'http://bit.ly', 'http://tinyurl', 'data:text/html',
            'javascript:void', 'vbscript:'
        ],
        'embedded_executables': [
            'mz\x90\x00', 'elf', '\x7felf', 'pe\x00\x00'
        ]
    }
    
    for threat_type, patterns in suspicious_patterns.items():
        for pattern in patterns:
            if pattern in text_content:
                threats.append(f"{threat_type}: {pattern}")
    
    return threats


async def analyze_embedded_content(content: bytes, filename: str) -> list:
    """Analyze document files for embedded suspicious content."""
    threats = []
    
    try:
        file_ext = Path(filename).suffix.lower()
        
        if file_ext == '.pdf':
            # Check for JavaScript in PDF
            text_content = content.decode('latin1', errors='ignore').lower()
            if '/js' in text_content or '/javascript' in text_content:
                threats.append("PDF contains JavaScript")
            if '/launch' in text_content:
                threats.append("PDF contains launch actions")
                
        elif file_ext in ['.docx', '.pptx']:
            # Check for macros or embedded objects in Office documents
            text_content = content.decode('latin1', errors='ignore').lower()
            if 'vbaproject' in text_content:
                threats.append("Document contains VBA macros")
            if 'oleobject' in text_content:
                threats.append("Document contains embedded objects")
                
    except Exception as e:
        logger.warning(f"Embedded content analysis failed for {filename}: {str(e)}")
    
    return threats


def get_expected_file_size_range(filename: str) -> tuple:
    """Get expected file size range based on file type."""
    file_ext = Path(filename).suffix.lower()
    
    # Size ranges in bytes (min, max)
    size_ranges = {
        '.txt': (1, 50 * 1024 * 1024),  # 1B to 50MB
        '.pdf': (100, 100 * 1024 * 1024),  # 100B to 100MB
        '.docx': (1000, 50 * 1024 * 1024),  # 1KB to 50MB
        '.pptx': (10000, 100 * 1024 * 1024),  # 10KB to 100MB
        '.mp4': (1000, 500 * 1024 * 1024),  # 1KB to 500MB
        '.mp3': (1000, 100 * 1024 * 1024),  # 1KB to 100MB
        '.wav': (1000, 200 * 1024 * 1024),  # 1KB to 200MB
    }
    
    return size_ranges.get(file_ext, (1, 100 * 1024 * 1024))  # Default: 1B to 100MB


def calculate_file_entropy(content: bytes) -> float:
    """Calculate Shannon entropy of file content."""
    if len(content) == 0:
        return 0
    
    # Count frequency of each byte
    frequency = {}
    for byte in content:
        frequency[byte] = frequency.get(byte, 0) + 1
    
    # Calculate entropy
    entropy = 0
    for count in frequency.values():
        probability = count / len(content)
        entropy -= probability * (probability).bit_length()
    
    return entropy


def is_filename_safe(filename: str) -> bool:
    """Check if filename contains safe patterns."""
    # Check for path traversal attempts
    if '..' in filename or '/' in filename or '\\' in filename:
        return False
    
    # Check for suspicious extensions
    suspicious_extensions = [
        '.exe', '.bat', '.cmd', '.com', '.scr', '.pif',
        '.vbs', '.js', '.jar', '.sh', '.ps1'
    ]
    
    filename_lower = filename.lower()
    for ext in suspicious_extensions:
        if filename_lower.endswith(ext):
            return False
    
    # Check for double extensions (file.pdf.exe)
    parts = filename.split('.')
    if len(parts) > 2:
        for i in range(len(parts) - 1):
            potential_ext = '.' + parts[i]
            if potential_ext in suspicious_extensions:
                return False
    
    return True

# Security Validation Endpoint
@app.post("/validate-file")
@limiter.limit("10/minute")
async def validate_file_security_endpoint(
    request: Request,
    file: UploadFile = File(...)
):
    """Validate file security and return detailed security report."""
    try:
        logger.info(f"🔍 Security validation requested for: {file.filename}")
        
        # Perform comprehensive security validation
        validation_results = await validate_file_security(file)
        
        # Add recommendations based on results
        recommendations = []
        if validation_results["security_score"] < 60:
            recommendations.append("File rejected due to security concerns")
        elif validation_results["security_score"] < 80:
            recommendations.append("File accepted with caution - monitor during processing")
        else:
            recommendations.append("File passed security validation")
        
        validation_results["recommendations"] = recommendations
        
        # Return detailed security report
        return {
            "status": "validation_complete",
            "filename": file.filename,
            "timestamp": datetime.now().isoformat(),
            **validation_results
        }
        
    except FileValidationError as e:
        raise e
    except Exception as e:
        logger.error(f"❌ Security validation endpoint error: {str(e)}")
        raise ContentProcessingError(f"Security validation failed: {str(e)}")


# Queue Management Endpoints
@app.get("/queue/stats")
@limiter.limit("30/minute")
async def get_queue_stats(request: Request):
    """Get processing queue statistics and status."""
    try:
        stats = processing_queue.get_queue_stats()
        
        return {
            "queue_stats": stats,
            "worker_count": processing_queue.max_workers,
            "queue_capacity": processing_queue.max_queue_size,
            "is_running": processing_queue.running,
            "timestamp": datetime.now().isoformat()
        }
        
    except Exception as e:
        logger.error(f"❌ Error getting queue stats: {str(e)}")
        raise ContentProcessingError(f"Failed to get queue stats: {str(e)}")


@app.get("/queue/jobs/active")
@limiter.limit("20/minute")
async def get_active_jobs(request: Request):
    """Get currently active processing jobs."""
    try:
        active_jobs = []
        
        for job_id, job in processing_queue.active_jobs.items():
            active_jobs.append({
                "job_id": job.job_id,
                "material_id": job.material_id,
                "priority": job.priority.name,
                "created_at": job.created_at.isoformat(),
                "retries": job.retries,
                "max_retries": job.max_retries
            })
        
        return {
            "active_jobs": active_jobs,
            "count": len(active_jobs),
            "timestamp": datetime.now().isoformat()
        }
        
    except Exception as e:
        logger.error(f"❌ Error getting active jobs: {str(e)}")
        raise ContentProcessingError(f"Failed to get active jobs: {str(e)}")


@app.post("/queue/priority/{job_id}")
@limiter.limit("5/minute")
async def update_job_priority(
    request: Request,
    job_id: str,
    priority: str = Form(...)
):
    """Update the priority of a queued job."""
    try:
        # Validate priority
        try:
            new_priority = JobPriority[priority.upper()]
        except KeyError:
            raise ContentProcessingError(
                f"Invalid priority '{priority}'. Valid options: LOW, NORMAL, HIGH, URGENT"
            )
        
        # Check if job is in any queue
        job_found = False
        for queue_priority, queue in processing_queue.queues.items():
            # This is a simple approach - in production you'd want a more efficient job lookup
            temp_jobs = []
            while not queue.empty():
                try:
                    job = await asyncio.wait_for(queue.get(), timeout=0.1)
                    if job.job_id == job_id:
                        job.priority = new_priority
                        job_found = True
                    temp_jobs.append(job)
                except asyncio.TimeoutError:
                    break
            
            # Put jobs back
            for job in temp_jobs:
                await processing_queue.add_job(job)
        
        if not job_found:
            raise ContentProcessingError(f"Job {job_id} not found in queue or already processing")
        
        return {
            "message": f"Job {job_id} priority updated to {new_priority.name}",
            "job_id": job_id,
            "new_priority": new_priority.name,
            "timestamp": datetime.now().isoformat()
        }
        
    except ContentProcessingError:
        raise
    except Exception as e:
        logger.error(f"❌ Error updating job priority: {str(e)}")
        raise ContentProcessingError(f"Failed to update job priority: {str(e)}")


# Duplicate Detection Endpoint
@app.post("/check-duplicates")
@limiter.limit("10/minute")
async def check_for_duplicates(
    request: Request,
    file: UploadFile = File(...),
    course_id: str = Form(None)
):
    """Check if uploaded file is a duplicate of existing content."""
    try:
        logger.info(f"🔍 Checking for duplicates: {file.filename}")
        
        # Read file content
        await file.seek(0)
        content = await file.read()
        await file.seek(0)
        
        # Extract text for comparison
        try:
            # Use a simplified text extraction for duplicate checking
            text_content = ""
            file_ext = Path(file.filename).suffix.lower()
            
            if file_ext == '.txt':
                text_content = content.decode('utf-8', errors='ignore')
            elif file_ext == '.pdf' and PDF_AVAILABLE:
                try:
                    import io
                    pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
                    text_parts = []
                    for page in pdf_reader.pages:
                        text_parts.append(page.extract_text())
                    text_content = '\n'.join(text_parts)
                except Exception as e:
                    logger.warning(f"PDF text extraction failed: {str(e)}")
                    text_content = ""
            else:
                # For other file types, use content as string representation
                text_content = content.decode('utf-8', errors='ignore')
            
        except Exception as e:
            logger.warning(f"Text extraction for duplicate check failed: {str(e)}")
            text_content = ""
        
        # Check for duplicates
        duplicates = await duplicate_detector.find_duplicates(
            content, text_content, file.filename, course_id
        )
        
        # Add file metadata
        duplicates["file_info"] = {
            "filename": file.filename,
            "size_bytes": len(content),
            "content_type": file.content_type,
            "file_hash": duplicate_detector.calculate_file_hash(content),
            "content_hash": duplicate_detector.calculate_content_hash(text_content) if text_content else None
        }
        
        return {
            "status": "duplicate_check_complete",
            "filename": file.filename,
            "course_id": course_id,
            "timestamp": datetime.now().isoformat(),
            **duplicates
        }
        
    except Exception as e:
        logger.error(f"❌ Duplicate check failed: {str(e)}")
        raise ContentProcessingError(f"Duplicate detection failed: {str(e)}")


@app.post("/process-with-duplicate-check")
@limiter.limit("3/minute")
async def process_material_with_duplicate_check(
    request: Request,
    material_id: str = Form(...),
    file_path: str = Form(...),
    processing_job_id: str = Form(...),
    force_processing: bool = Form(False)
):
    """Process material with automatic duplicate detection."""
    try:
        logger.info(f"🔄 Processing with duplicate check: {material_id}")
        
        # Validate that the file exists
        file_path_obj = Path(file_path)
        if not file_path_obj.exists():
            raise FileValidationError(f"File not found at path: {file_path}")
        
        # Read file for duplicate checking
        with open(file_path, 'rb') as f:
            content = f.read()
        
        # Extract text for duplicate checking
        try:
            text_content = await extract_text_from_file(file_path)
        except Exception as e:
            logger.warning(f"Text extraction failed for duplicate check: {str(e)}")
            text_content = ""
        
        # Check for duplicates unless forced
        if not force_processing:
            duplicates = await duplicate_detector.find_duplicates(
                content, text_content, file_path_obj.name
            )
            
            if duplicates["is_duplicate"]:
                return {
                    "status": "duplicate_detected",
                    "message": "Duplicate content detected - processing skipped",
                    "material_id": material_id,
                    "processing_job_id": processing_job_id,
                    "duplicates": duplicates,
                    "force_processing_available": True
                }
        
        # If no duplicates or forced processing, proceed with normal processing
        # (This would call the existing process_existing_material logic)
        
        # Verify material exists in database
        try:
            conn = get_db_connection()
            cursor = conn.cursor()
            cursor.execute("""
                SELECT cm.title, cm.material_type, cm.course_id
                FROM course_material cm
                WHERE cm.id = %s
            """, (material_id,))
            
            result = cursor.fetchone()
            cursor.close()
            conn.close()
            
            if not result:
                raise ContentProcessingError(
                    f"Course material {material_id} not found in database",
                    error_code="MATERIAL_NOT_FOUND"
                )
            
            # Handle both tuple and dict results from cursor
            if isinstance(result, dict):
                material_title = result['title']
                material_type = result['material_type']
                course_id = result['course_id']
            else:
                material_title, material_type, course_id = result
                
            logger.info(f"✅ Material found: {material_title} ({material_type})")
            
        except ContentProcessingError:
            raise
        except Exception as e:
            logger.error(f"❌ Database verification failed: {str(e)}")
            raise DatabaseError(f"Failed to verify material: {str(e)}")
        
        # Add job to processing queue
        try:
            # Determine priority based on material type
            priority = JobPriority.NORMAL
            if material_type in ['exam', 'assignment']:
                priority = JobPriority.HIGH
            elif material_type == 'syllabus':
                priority = JobPriority.URGENT
            
            # Create processing job
            job = ProcessingJob(
                job_id=processing_job_id,
                material_id=material_id,
                file_path=file_path,
                priority=priority,
                created_at=datetime.now()
            )
            
            # Add to queue
            success = await processing_queue.add_job(job)
            if not success:
                raise AIProcessingError("Processing queue is full, please try again later")
            
            logger.info(f"🚀 Material {material_id} queued for processing with {priority.name} priority")
        except Exception as e:
            logger.error(f"❌ Failed to queue processing job: {str(e)}")
            raise AIProcessingError(f"Failed to queue AI processing: {str(e)}")
        
        return {
            "status": "queued_for_processing",
            "message": "Material queued for processing successfully",
            "material_id": material_id,
            "processing_job_id": processing_job_id,
            "priority": priority.name,
            "material_title": material_title,
            "duplicate_check_passed": True
        }
        
    except (FileValidationError, ContentProcessingError, DatabaseError, AIProcessingError) as e:
        raise e
    except Exception as e:
        logger.error(f"❌ Unexpected error in process with duplicate check: {str(e)}")
        raise ContentProcessingError(f"Processing with duplicate check failed: {str(e)}")


# Content Processing Endpoints with Enhanced Error Handling

@app.post("/process-content")
@limiter.limit("5/minute")
async def process_content_legacy(
    request: Request,
    title: str = Form(...),
    course_id: Optional[int] = Form(None),
    week_number: Optional[int] = Form(None),
    material_type: Optional[str] = Form(None),
    file: UploadFile = File(...)
):
    """Legacy endpoint for direct file uploads - will be deprecated"""
    # This is the existing implementation for backward compatibility
    # TODO: Remove once Next.js migration is complete
    """Process uploaded content with enhanced validation and error handling."""
    try:
        logger.info(f"🔄 Processing content upload: {title}")
        
        # Enhanced file validation
        await validate_file_enhanced(file)
        
        # Validate request data
        request_data = ContentUploadRequest(
            title=title,
            course_id=course_id,
            week_number=week_number,
            material_type=material_type
        )
        
        # Save file with secure filename
        safe_filename = secure_filename(file.filename)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        unique_filename = f"{timestamp}_{safe_filename}"
        file_path = Config.UPLOAD_DIR / unique_filename
        
        # Save uploaded file
        with open(file_path, "wb") as buffer:
            content = await file.read()
            buffer.write(content)
        
        logger.info(f"✅ File saved: {file_path}")
        
        # Insert into database with enhanced error handling
        try:
            conn = get_db_connection()
            cursor = conn.cursor()
            cursor.execute("""
                INSERT INTO course_material (title, course_id, week_number, material_type, content_url, uploaded_by_id, created_at)
                VALUES (%s, %s, %s, %s, %s, %s, %s) RETURNING id
            """, (
                request_data.title,
                request_data.course_id,
                request_data.week_number,
                request_data.material_type or 'lecture',
                str(file_path),
                'system',  # We'll need to get actual user ID later
                datetime.now()
            ))
            result = cursor.fetchone()
            if result is None:
                raise DatabaseError("No ID returned from database insertion")
            
            # Handle RealDictCursor result (returns dict) vs regular cursor result (returns tuple)
            if isinstance(result, dict):
                material_id = result['id']
            else:
                material_id = result[0]
                
            if not material_id:
                raise DatabaseError("Invalid material ID returned from database")
                
            conn.commit()
            cursor.close()
            conn.close()
            
            logger.info(f"✅ Material inserted with ID: {material_id}")
            
        except Exception as e:
            logger.error(f"❌ Database insertion failed: {str(e)}")
            # Clean up uploaded file on database failure
            if file_path.exists():
                file_path.unlink()
            raise DatabaseError(f"Failed to save material metadata: {str(e)}")
        
        # Create processing job
        job_type_mapping = {
            'pdf': 'pdf_processing',
            'video': 'video_transcription', 
            'audio': 'video_transcription',
            'interactive': 'interactive_parsing',
            'text': 'pdf_processing'
        }
        
        job_type = job_type_mapping.get(request_data.material_type, 'pdf_processing')
        
        try:
            conn = get_db_connection()
            cursor = conn.cursor()
            cursor.execute("""
                INSERT INTO ai_processing_job (course_material_id, job_type, status, created_at)
                VALUES (%s, %s, 'pending', %s) RETURNING id
            """, (material_id, job_type, datetime.now()))
            
            result = cursor.fetchone()
            job_id = result['id'] if isinstance(result, dict) else result[0]
            conn.commit()
            cursor.close()
            conn.close()
            
            logger.info(f"✅ Processing job created with ID: {job_id}")
            
        except Exception as e:
            logger.error(f"❌ Job creation failed: {str(e)}")
            raise DatabaseError(f"Failed to create processing job: {str(e)}")
        
        # Start background processing
        try:
            await start_background_processing(material_id, job_id, str(file_path))
            logger.info(f"🚀 Background processing started for job {job_id}")
        except Exception as e:
            logger.error(f"❌ Background processing failed: {str(e)}")
            raise AIProcessingError(f"Failed to start AI processing: {str(e)}")
        
        return {
            "message": "Content uploaded and processing started",
            "material_id": material_id,
            "processing_id": job_id,
            "status": "processing"
        }
        
    except (FileValidationError, ContentProcessingError, DatabaseError, AIProcessingError) as e:
        # These are our custom exceptions, re-raise them
        raise e
    except Exception as e:
        logger.error(f"❌ Unexpected error in process_content: {str(e)}")
        raise ContentProcessingError(f"Unexpected error during processing: {str(e)}")

async def process_content(file_path: str, material_id: str, processing_job_id: str) -> dict:
    """
    Core content processing function that handles file analysis and AI processing.
    
    Args:
        file_path: Path to the file to process
        material_id: ID of the course material
        processing_job_id: ID of the processing job
        
    Returns:
        dict: Processing results including extracted text, summary, and concepts
    """
    ai_processed_id = None
    
    try:
        logger.info(f"🚀 Starting process_content for material {material_id}")
        logger.info(f"📁 File path: {file_path}")
        logger.info(f"🆔 Processing job ID: {processing_job_id}")
        
        # Get AI stack instance
        logger.info("🔄 Getting AI stack instance...")
        ai_stack = await get_ai_stack()
        logger.info("✅ AI stack instance obtained")
        
        # Extract text from file
        logger.info(f"📄 Extracting text from file: {file_path}")
        logger.info(f"📁 File exists check: {os.path.exists(file_path)}")
        if os.path.exists(file_path):
            logger.info(f"📊 File size: {os.path.getsize(file_path)} bytes")
        
        extracted_text = await extract_text_from_file(file_path)
        
        if not extracted_text.strip():
            logger.error("❌ No text could be extracted from the file")
            raise ContentProcessingError("No text could be extracted from the file")
        
        logger.info(f"✅ Text extraction completed ({len(extracted_text)} characters)")
        logger.info(f"📝 First 200 chars: {extracted_text[:200]}...")
        
        # Generate AI analysis (summary and key concepts)
        logger.info("🧠 Generating AI analysis...")
        try:
            analysis_result = await ai_stack.analyze_content(extracted_text[:4000], "educational")
            logger.info(f"🧠 Analysis result: {analysis_result}")
        except Exception as e:
            logger.error(f"❌ AI analysis failed: {e}")
            raise AIProcessingError(f"AI analysis failed: {str(e)}")
        
        # Extract summary and concepts from analysis
        ai_summary = analysis_result.get('response', '') if analysis_result.get('success') else ''
        logger.info(f"📋 AI summary length: {len(ai_summary)} characters")
        
        # Parse key concepts from the analysis response
        logger.info("🔍 Extracting key concepts from analysis...")
        concepts_prompt = f"""Extract key concepts from this educational content as a simple comma-separated list:

{extracted_text[:2000]}

Return only the key concepts as a comma-separated list, nothing else."""
        
        try:
            concepts_result = await ai_stack.generate_llm_response(concepts_prompt)
            logger.info(f"🔍 Concepts result: {concepts_result}")
            key_concepts = concepts_result.get('response', '').split(',') if concepts_result.get('success') else []
            key_concepts = [concept.strip() for concept in key_concepts if concept.strip()]
            logger.info(f"🏷️ Extracted {len(key_concepts)} key concepts: {key_concepts}")
        except Exception as e:
            logger.error(f"❌ Key concepts extraction failed: {e}")
            key_concepts = []
        
        # Generate embeddings
        logger.info("🔗 Generating embeddings...")
        try:
            embeddings_response = await ai_stack.generate_embeddings(extracted_text[:8000])
            logger.info(f"🔗 Embeddings response type: {type(embeddings_response)}")
            logger.info(f"🔗 Embeddings response keys: {embeddings_response.keys() if isinstance(embeddings_response, dict) else 'Not a dict'}")
        except Exception as e:
            logger.error(f"❌ Embeddings generation failed: {e}")
            embeddings_response = {'success': False, 'error': str(e)}
        
        # Process embeddings response
        if isinstance(embeddings_response, dict) and 'embedding' in embeddings_response:
            embeddings_array = embeddings_response['embedding']
        elif isinstance(embeddings_response, list):
            embeddings_array = embeddings_response
        else:
            logger.warning(f"⚠️ Unexpected embeddings format: {type(embeddings_response)}")
            embeddings_array = []
        
        # Save processed content to database
        logger.info("💾 Saving processed content to database...")
        try:
            conn = get_db_connection()
            logger.info("✅ Database connection established")
            cursor = conn.cursor()
            logger.info("✅ Database cursor created")
        except Exception as e:
            logger.error(f"❌ Database connection failed: {e}")
            raise DatabaseError(f"Database connection failed: {str(e)}")
        
        try:
            # Save processed content
            logger.info("💾 Inserting AI processed content...")
            cursor.execute("""
                INSERT INTO ai_processed_content 
                (course_material_id, extracted_text, ai_summary, key_concepts, created_at)
                VALUES (%s, %s, %s, %s, %s) RETURNING id
            """, (
                material_id,
                extracted_text,
                ai_summary,
                json.dumps(key_concepts),
                datetime.now()
            ))
            
            result = cursor.fetchone()
            ai_processed_id = result['id'] if isinstance(result, dict) else result[0]
            logger.info(f"✅ AI processed content saved with ID: {ai_processed_id}")
            
            # Save embeddings if they exist
            if embeddings_array and len(embeddings_array) == 768:
                logger.info("💾 Inserting content embeddings...")
                # Convert embeddings to vector format for pgvector
                vector_str = '[' + ','.join(map(str, embeddings_array)) + ']'
                
                cursor.execute("""
                    INSERT INTO content_embedding 
                    (course_material_id, ai_processed_id, chunk_text, chunk_index, chunk_type, embedding, created_at)
                    VALUES (%s, %s, %s, %s, %s, %s, %s)
                """, (
                    material_id, 
                    ai_processed_id, 
                    extracted_text[:1000], 
                    1, 
                    'content', 
                    vector_str, 
                    datetime.now()
                ))
                logger.info(f"✅ Embeddings saved as vector ({len(embeddings_array)} dimensions)")
            else:
                logger.warning(f"⚠️ Skipping embeddings: invalid format or dimensions ({len(embeddings_array) if embeddings_array else 0})")
            
            conn.commit()
            logger.info(f"✅ Content processing completed successfully (AI processed ID: {ai_processed_id})")
            logger.info(f"📊 Final stats: Text={len(extracted_text)} chars, Summary={len(ai_summary)} chars, Concepts={len(key_concepts)}")
            
            return {
                "extracted_text": extracted_text[:500] + "..." if len(extracted_text) > 500 else extracted_text,
                "ai_summary": ai_summary,
                "key_concepts": key_concepts,
                "embeddings_count": len(embeddings_array) if embeddings_array else 0,
                "ai_processed_id": ai_processed_id,
                "status": "completed"
            }
            
        except Exception as db_error:
            conn.rollback()
            raise DatabaseError(f"Database save failed: {str(db_error)}")
        finally:
            cursor.close()
            conn.close()
        
    except Exception as e:
        logger.error(f"❌ Content processing failed: {str(e)}")
        logger.error(f"❌ Exception type: {type(e).__name__}")
        logger.error(f"❌ Exception details: {str(e)}")
        import traceback
        logger.error(f"❌ Stack trace: {traceback.format_exc()}")
        raise AIProcessingError(f"Processing failed: {str(e)}")

@app.post("/process-material")
@limiter.limit("5/minute")
async def process_s3_material(
    request: Request,
    material_id: str = Form(...),
    s3_key: str = Form(...),
    s3_bucket: str = Form(...),
    file_type: str = Form(...),
    processing_job_id: str = Form(...)
):
    """Process course material from S3 storage with AI analysis and embedding generation."""
    temp_file_path = None
    
    try:
        logger.info(f"🔄 Processing S3 material: {material_id} from s3://{s3_bucket}/{s3_key}")
        
        # Verify S3 service is available
        if not s3_service.health_check():
            raise HTTPException(
                status_code=503,
                detail="S3 service unavailable"
            )
        
        # Check if file exists in S3
        if not s3_service.file_exists(s3_key):
            raise FileNotFoundError(f"File not found in S3: s3://{s3_bucket}/{s3_key}")
        
        # Update processing job status to "processing"
        try:
            conn = get_db_connection()
            cursor = conn.cursor()
            cursor.execute("""
                UPDATE ai_processing_job
                SET status = 'processing', started_at = CURRENT_TIMESTAMP
                WHERE id = %s
            """, (processing_job_id,))
            conn.commit()
            cursor.close()
            conn.close()
        except Exception as e:
            logger.warning(f"Failed to update job status to processing: {e}")
        
        # Download file from S3 to temporary location
        logger.info(f"📥 Downloading file from S3: {s3_key}")
        temp_file_path, original_filename = s3_service.download_file_to_temp(s3_key)
        
        logger.info(f"✅ File downloaded successfully: {temp_file_path}")
        
        # Verify material exists in database
        try:
            conn = get_db_connection()
            cursor = conn.cursor()
            cursor.execute("""
                SELECT cm.title, cm.material_type, cm.course_id
                FROM course_material cm
                WHERE cm.id = %s
            """, (material_id,))
            
            result = cursor.fetchone()
            cursor.close()
            conn.close()
            
            if not result:
                raise ContentProcessingError(
                    f"Course material {material_id} not found in database",
                    error_code="MATERIAL_NOT_FOUND"
                )
            
            # Handle both tuple and dict results from cursor
            if isinstance(result, dict):
                material_title = result['title']
                material_type = result['material_type']
                course_id = result['course_id']
            else:
                material_title, material_type, course_id = result
                
            logger.info(f"✅ Material found: {material_title} ({material_type})")
            
        except ContentProcessingError:
            raise
        except Exception as e:
            logger.error(f"❌ Database verification failed: {str(e)}")
            raise DatabaseError(f"Failed to verify material: {str(e)}")
        
        # Process the downloaded file directly (simplified approach)
        try:
            # Determine priority based on material type
            priority = JobPriority.NORMAL
            if material_type in ['exam', 'assignment']:
                priority = JobPriority.HIGH
            elif material_type == 'syllabus':
                priority = JobPriority.URGENT
            
            # Process the file using the temporary path
            logger.info(f"🧠 Starting AI processing for {material_title}")
            
            # Process content using existing process_content function
            processing_result = await process_content(
                temp_file_path,
                material_id,
                processing_job_id
            )
            
            # Update job status to completed
            try:
                conn = get_db_connection()
                cursor = conn.cursor()
                cursor.execute("""
                    UPDATE ai_processing_job
                    SET status = 'completed', completed_at = CURRENT_TIMESTAMP
                    WHERE id = %s
                """, (processing_job_id,))
                conn.commit()
                cursor.close()
                conn.close()
                logger.info(f"✅ Processing job {processing_job_id} marked as completed")
            except Exception as e:
                logger.warning(f"Failed to update job status to completed: {e}")
            
            logger.info(f"🎉 Successfully processed S3 material {material_id}")
            
            # Return successful processing result
            return {
                "status": "completed",
                "message": "Material processed successfully",
                "material_id": material_id,
                "processing_job_id": processing_job_id,
                "material_title": material_title,
                "processing_result": processing_result
            }
            
        except Exception as e:
            logger.error(f"❌ Failed to process S3 material: {str(e)}")
            
            # Update job status to failed
            try:
                conn = get_db_connection()
                cursor = conn.cursor()
                cursor.execute("""
                    UPDATE ai_processing_job
                    SET status = 'failed', completed_at = CURRENT_TIMESTAMP, error_message = %s
                    WHERE id = %s
                """, (str(e), processing_job_id))
                conn.commit()
                cursor.close()
                conn.close()
            except Exception as db_e:
                logger.warning(f"Failed to update job status to failed: {db_e}")
            
            raise AIProcessingError(f"Failed to process S3 material: {str(e)}")
        
    except (FileValidationError, ContentProcessingError, DatabaseError, AIProcessingError) as e:
        raise e
    except Exception as e:
        logger.error(f"❌ Unexpected error in process_s3_material: {str(e)}")
        raise ContentProcessingError(f"Unexpected error during S3 processing: {str(e)}")
    finally:
        # Always cleanup temporary file
        if temp_file_path and os.path.exists(temp_file_path):
            try:
                s3_service.cleanup_temp_file(temp_file_path)
                logger.info(f"🧹 Cleaned up temporary file: {temp_file_path}")
            except Exception as cleanup_e:
                logger.warning(f"Failed to cleanup temp file: {cleanup_e}")

@app.get("/processing-status/{processing_id}")
@limiter.limit("30/minute")
async def get_processing_status(request: Request, processing_id: str):
    """Get processing status with enhanced error handling."""
    try:
        logger.info(f"📊 Checking status for processing ID: {processing_id}")
        
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute("""
            SELECT aj.status, aj.started_at, aj.completed_at, aj.error_message,
                   cm.title, cm.id as material_id
            FROM ai_processing_job aj
            JOIN course_material cm ON aj.course_material_id = cm.id
            WHERE aj.id = %s
        """, (processing_id,))
        
        result = cursor.fetchone()
        cursor.close()
        conn.close()
        
        if not result:
            raise ContentProcessingError(
                f"Processing job {processing_id} not found",
                error_code="JOB_NOT_FOUND"
            )
        
        # Handle both tuple and dict results from cursor
        if isinstance(result, dict):
            status = result['status']
            started_at = result['started_at']
            completed_at = result['completed_at']
            error_message = result['error_message']
            title = result['title']
            material_id = result['material_id']
        else:
            status, started_at, completed_at, error_message, title, material_id = result
        
        # Helper function to safely format datetime
        def safe_isoformat(dt_value):
            if dt_value is None:
                return None
            if hasattr(dt_value, 'isoformat'):
                return dt_value.isoformat()
            return str(dt_value)  # Return as string if already a string
        
        response_data = {
            "processing_id": processing_id,
            "material_id": material_id,
            "title": title,
            "status": status,
            "started_at": safe_isoformat(started_at),
            "completed_at": safe_isoformat(completed_at)
        }
        
        if error_message:
            response_data["error_message"] = error_message
            
        logger.info(f"✅ Status retrieved for job {processing_id}: {status}")
        return response_data
        
    except ContentProcessingError as e:
        raise e
    except Exception as e:
        logger.error(f"❌ Error getting processing status: {str(e)}")
        raise DatabaseError(f"Failed to retrieve processing status: {str(e)}")

@app.post("/search")
@limiter.limit("20/minute") 
async def search_content(request: Request, search_request: SearchRequest):
    """Enhanced semantic search with validation."""
    try:
        logger.info(f"🔍 Searching for: {search_request.query}")
        
        if not search_request.query.strip():
            raise ContentProcessingError(
                "Search query cannot be empty",
                error_code="INVALID_QUERY"
            )
        
        # Generate query embedding
        try:
            query_embedding_response = await ai_stack.generate_embeddings(search_request.query)
            if not query_embedding_response:
                raise AIProcessingError("Failed to generate query embeddings")
            
            # Extract embedding array from AI response
            if isinstance(query_embedding_response, dict) and 'embedding' in query_embedding_response:
                query_embedding = query_embedding_response['embedding']
            elif isinstance(query_embedding_response, list):
                query_embedding = query_embedding_response
            else:
                raise AIProcessingError("Invalid embedding format from AI response")
                
        except Exception as e:
            logger.error(f"❌ Embedding generation failed: {str(e)}")
            raise AIProcessingError(f"Search embedding failed: {str(e)}")
        
        # Perform semantic search
        try:
            conn = get_db_connection()
            cursor = conn.cursor()
            cursor.execute("""
                SELECT cm.id, cm.title::TEXT, cm.material_type, cm.week_number,
                       apc.ai_summary, apc.key_concepts,
                       (ce.embedding <=> %s::vector) as similarity_score
                FROM course_material cm
                JOIN ai_processed_content apc ON cm.id = apc.course_material_id
                JOIN content_embedding ce ON apc.id = ce.ai_processed_id
                JOIN course c ON cm.course_id = c.id
                WHERE (%s IS NULL OR c.course_code = %s)
                ORDER BY ce.embedding <=> %s::vector
                LIMIT %s
            """, (
                query_embedding,
                search_request.course_filter,
                search_request.course_filter,
                query_embedding,
                search_request.limit
            ))
            
            results = cursor.fetchall()
            cursor.close()
            conn.close()
            
            search_results = []
            for row in results:
                # Handle both tuple and dict results from cursor
                if isinstance(row, dict):
                    material_id = row['id']
                    title = row['title']
                    material_type = row['material_type']
                    week_number = row['week_number']
                    ai_summary = row['ai_summary']
                    key_concepts = row['key_concepts']
                    similarity_score = row['similarity_score']
                else:
                    material_id, title, material_type, week_number, ai_summary, key_concepts, similarity_score = row
                
                search_results.append({
                    "material_id": material_id,
                    "title": title,
                    "material_type": material_type,
                    "week_number": week_number,
                    "ai_summary": ai_summary,
                    "key_concepts": key_concepts,
                    "similarity_score": float(similarity_score),
                    "relevance": "high" if similarity_score < 0.3 else "medium" if similarity_score < 0.6 else "low"
                })
            
            logger.info(f"✅ Search completed: {len(search_results)} results found")
            
            return {
                "query": search_request.query,
                "results_count": len(search_results),
                "results": search_results
            }
            
        except Exception as e:
            logger.error(f"❌ Search query failed: {str(e)}")
            raise DatabaseError(f"Search operation failed: {str(e)}")
        
    except (ContentProcessingError, AIProcessingError, DatabaseError) as e:
        raise e
    except Exception as e:
        logger.error(f"❌ Unexpected search error: {str(e)}")
        raise ContentProcessingError(f"Search failed: {str(e)}")

@app.get("/courses/{course_code}/materials")
@limiter.limit("30/minute")
async def get_course_materials(request: Request, course_code: str):
    """Get course materials with enhanced error handling."""
    try:
        logger.info(f"📚 Getting materials for course: {course_code}")
        
        if not course_code.strip():
            raise ContentProcessingError(
                "Course code cannot be empty",
                error_code="INVALID_COURSE_CODE"
            )
        
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute("""
            SELECT cm.id, cm.title, cm.material_type, cm.week_number, 
                   cm.created_at, apc.ai_summary, apc.key_concepts
            FROM course_material cm
            LEFT JOIN ai_processed_content apc ON cm.id = apc.course_material_id
            WHERE cm.course_id = (
                SELECT id FROM course WHERE course_code = %s LIMIT 1
            )
            ORDER BY cm.week_number, cm.created_at
        """, (course_code,))
        
        results = cursor.fetchall()
        cursor.close()
        conn.close()
        
        if not results:
            logger.warning(f"⚠️ No materials found for course: {course_code}")
            return {
                "course_code": course_code,
                "materials_count": 0,
                "materials": []
            }
        
        materials = []
        for row in results:
            material_id, title, material_type, week_number, created_at, ai_summary, key_concepts = row
            materials.append({
                "material_id": material_id,
                "title": title,
                "material_type": material_type,
                "week_number": week_number,
                "created_at": created_at.isoformat() if created_at else None,
                "ai_summary": ai_summary,
                "key_concepts": key_concepts,
                "has_ai_processing": ai_summary is not None
            })
        
        logger.info(f"✅ Retrieved {len(materials)} materials for course {course_code}")
        
        return {
            "course_code": course_code,
            "materials_count": len(materials),
            "materials": materials
        }
        
    except ContentProcessingError as e:
        raise e
    except Exception as e:
        logger.error(f"❌ Error getting course materials: {str(e)}")
        raise DatabaseError(f"Failed to retrieve course materials: {str(e)}")


@app.post("/retry-processing/{job_id}")
@limiter.limit("3/minute")
async def retry_failed_processing(request: Request, job_id: str):
    """Retry a failed processing job with enhanced recovery capabilities."""
    try:
        logger.info(f"🔄 Retrying failed processing job: {job_id}")
        
        # Get job details and verify it can be retried
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute("""
            SELECT aj.status, aj.error_message, aj.metadata, 
                   cm.id as material_id, cm.file_path
            FROM ai_processing_job aj
            JOIN course_material cm ON aj.course_material_id = cm.id
            WHERE aj.id = %s
        """, (job_id,))
        
        result = cursor.fetchone()
        cursor.close()
        conn.close()
        
        if not result:
            raise ContentProcessingError(
                f"Processing job {job_id} not found",
                error_code="JOB_NOT_FOUND"
            )
        
        # Handle both tuple and dict results
        if isinstance(result, dict):
            status = result['status']
            error_message = result['error_message']
            metadata = result['metadata']
            material_id = result['material_id']
            file_path = result['file_path']
        else:
            status, error_message, metadata, material_id, file_path = result
        
        # Check if job can be retried
        if status != 'failed':
            raise ContentProcessingError(
                f"Job {job_id} is not in failed state (current: {status})",
                error_code="INVALID_STATUS"
            )
        
        # Parse metadata to check retry eligibility
        can_retry = True
        if metadata:
            try:
                metadata_dict = json.loads(metadata) if isinstance(metadata, str) else metadata
                can_retry = metadata_dict.get('can_retry', True)
            except (json.JSONDecodeError, AttributeError):
                pass
        
        if not can_retry:
            raise ContentProcessingError(
                f"Job {job_id} cannot be retried due to permanent failure",
                error_code="RETRY_NOT_ALLOWED"
            )
        
        # Verify file still exists
        if not Path(file_path).exists():
            raise FileValidationError(f"Source file no longer exists: {file_path}")
        
        # Clean up any existing failed processing data
        await cleanup_failed_job_data(material_id)
        
        # Reset job status and start retry
        try:
            conn = get_db_connection()
            cursor = conn.cursor()
            cursor.execute("""
                UPDATE ai_processing_job 
                SET status = 'pending', 
                    started_at = NULL, 
                    completed_at = NULL,
                    error_message = NULL,
                    metadata = COALESCE(metadata, '{}'::jsonb) || %s::jsonb
                WHERE id = %s
            """, (
                json.dumps({
                    "retry_attempt": True,
                    "retry_timestamp": datetime.now().isoformat(),
                    "original_error": error_message
                }),
                job_id
            ))
            conn.commit()
            cursor.close()
            conn.close()
            
            logger.info(f"✅ Job {job_id} reset for retry")
            
        except Exception as db_error:
            raise DatabaseError(f"Failed to reset job for retry: {str(db_error)}")
        
        # Start background processing
        try:
            await start_background_processing(material_id, job_id, file_path)
            logger.info(f"🚀 Retry processing started for job {job_id}")
        except Exception as e:
            logger.error(f"❌ Retry processing failed: {str(e)}")
            raise AIProcessingError(f"Failed to start retry processing: {str(e)}")
        
        return {
            "message": "Processing job retry started successfully",
            "job_id": job_id,
            "material_id": material_id,
            "status": "processing",
            "retry_attempt": True
        }
        
    except (FileValidationError, ContentProcessingError, DatabaseError, AIProcessingError) as e:
        raise e
    except Exception as e:
        logger.error(f"❌ Unexpected error in retry processing: {str(e)}")
        raise ContentProcessingError(f"Retry operation failed: {str(e)}")


async def cleanup_failed_job_data(material_id: str):
    """Clean up any partial data from failed processing attempts."""
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        
        # Remove any partial AI processed content
        cursor.execute("""
            DELETE FROM content_embedding 
            WHERE course_material_id = %s
        """, (material_id,))
        
        cursor.execute("""
            DELETE FROM ai_processed_content 
            WHERE course_material_id = %s
        """, (material_id,))
        
        conn.commit()
        cursor.close()
        conn.close()
        
        logger.info(f"🧹 Cleaned up partial data for material: {material_id}")
        
    except Exception as e:
        logger.warning(f"⚠️ Failed to cleanup partial data for {material_id}: {str(e)}")


@app.get("/processing-status/{job_id}/stream")
async def stream_processing_status(job_id: str):
    """Stream real-time processing status updates via Server-Sent Events"""
    async def event_generator():
        # Create a queue for this client
        status_queue = asyncio.Queue(maxsize=100)
        
        try:
            # Subscribe to status updates
            status_tracker.subscribe(job_id, status_queue)
            
            # Send current status first
            current_status = status_tracker.get_current_status(job_id)
            if current_status:
                yield {
                    "event": "status_update",
                    "data": json.dumps(current_status)
                }
            else:
                # If no status available, send initial pending status
                yield {
                    "event": "status_update", 
                    "data": json.dumps({
                        "status": "pending",
                        "message": "Waiting for processing to start",
                        "timestamp": datetime.now().isoformat()
                    })
                }
            
            # Stream updates as they come
            while True:
                try:
                    # Wait for status update with timeout
                    status_update = await asyncio.wait_for(
                        status_queue.get(), 
                        timeout=30.0  # 30 second timeout
                    )
                    
                    yield {
                        "event": "status_update",
                        "data": json.dumps(status_update)
                    }
                    
                    # If processing is complete or failed, send one more update and close
                    if status_update.get("status") in ["completed", "failed"]:
                        await asyncio.sleep(1)  # Give client time to process
                        break
                        
                except asyncio.TimeoutError:
                    # Send heartbeat to keep connection alive
                    yield {
                        "event": "heartbeat",
                        "data": json.dumps({
                            "timestamp": datetime.now().isoformat()
                        })
                    }
                    
        except asyncio.CancelledError:
            logger.info(f"Client disconnected from job {job_id} status stream")
        except Exception as e:
            logger.error(f"Error in status stream for job {job_id}: {str(e)}")
            yield {
                "event": "error",
                "data": json.dumps({
                    "error": "Stream error occurred",
                    "timestamp": datetime.now().isoformat()
                })
            }
        finally:
            # Clean up subscription
            status_tracker.unsubscribe(job_id, status_queue)
    
    return EventSourceResponse(event_generator())


# Helper function for background processing
async def start_background_processing(material_id: str, job_id: str, file_path: str):
    """Start background AI processing with comprehensive error handling and recovery."""
    conn = None
    ai_processed_id = None
    processing_stage = "initialization"
    
    try:
        # Update job status to processing
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute("""
            UPDATE ai_processing_job 
            SET status = 'processing', started_at = %s 
            WHERE id = %s
        """, (datetime.now(), job_id))
        conn.commit()
        cursor.close()
        
        # Send real-time status update
        status_tracker.update_status(job_id, {
            "status": "processing",
            "stage": "initialization",
            "message": "Starting content processing...",
            "progress": 0
        })
        
        # Stage 1: Text Extraction with retry
        processing_stage = "text_extraction"
        status_tracker.update_status(job_id, {
            "status": "processing",
            "stage": "text_extraction",
            "message": "Extracting text from file...",
            "progress": 20
        })
        extracted_text = await retry_with_backoff(
            lambda: extract_text_from_file(file_path),
            max_retries=3,
            stage="text_extraction"
        )
        logger.info(f"✅ Text extraction completed: {len(extracted_text)} characters")
        
        # Stage 2: AI Summary Generation with retry
        processing_stage = "ai_summary"
        status_tracker.update_status(job_id, {
            "status": "processing",
            "stage": "ai_summary",
            "message": "Generating AI summary...",
            "progress": 40
        })
        ai_summary = await retry_with_backoff(
            lambda: ai_stack.generate_llm_response(
                f"Summarize this educational content in 2-3 sentences:\n\n{extracted_text[:2000]}"
            ),
            max_retries=3,
            stage="ai_summary"
        )
        logger.info("✅ AI summary generation completed")
        
        # Stage 3: Key Concepts Analysis with retry
        processing_stage = "key_concepts"
        status_tracker.update_status(job_id, {
            "status": "processing",
            "stage": "key_concepts",
            "message": "Analyzing key concepts...",
            "progress": 60
        })
        key_concepts = await retry_with_backoff(
            lambda: ai_stack.analyze_content(extracted_text[:3000]),
            max_retries=3,
            stage="key_concepts"
        )
        logger.info("✅ Key concepts analysis completed")
        
        # Stage 4: Embeddings Generation with retry
        processing_stage = "embeddings"
        status_tracker.update_status(job_id, {
            "status": "processing",
            "stage": "embeddings",
            "message": "Generating embeddings for semantic search...",
            "progress": 80
        })
        embeddings_response = await retry_with_backoff(
            lambda: ai_stack.generate_embeddings(extracted_text[:8000]),
            max_retries=3,
            stage="embeddings"
        )
        
        # Extract embedding array from AI response
        if isinstance(embeddings_response, dict) and 'embedding' in embeddings_response:
            embeddings_array = embeddings_response['embedding']
        elif isinstance(embeddings_response, list):
            embeddings_array = embeddings_response
        else:
            logger.warning(f"⚠️ Unexpected embeddings format: {type(embeddings_response)}, using fallback")
            embeddings_array = []
        
        logger.info("✅ Embeddings generation completed")
        
        # Stage 5: Database Operations with transaction safety
        processing_stage = "database_save"
        status_tracker.update_status(job_id, {
            "status": "processing",
            "stage": "database_save",
            "message": "Saving processed content to database...",
            "progress": 95
        })
        
        # Serialize key_concepts for database storage
        if isinstance(key_concepts, dict):
            key_concepts_array = [f"{k}: {v}" for k, v in key_concepts.items()]
        elif isinstance(key_concepts, list):
            key_concepts_array = [str(item) for item in key_concepts]
        else:
            key_concepts_array = [str(key_concepts)] if key_concepts else []
        
        # Use database transaction for atomicity
        conn.autocommit = False
        cursor = conn.cursor()
        
        try:
            # Save processed content
            cursor.execute("""
                INSERT INTO ai_processed_content 
                (course_material_id, extracted_text, ai_summary, key_concepts, created_at)
                VALUES (%s, %s, %s, %s, %s) RETURNING id
            """, (
                material_id,
                extracted_text,
                ai_summary.get('response', '') if ai_summary else '',
                key_concepts_array,
                datetime.now()
            ))
            
            result = cursor.fetchone()
            ai_processed_id = result['id'] if isinstance(result, dict) else result[0]
            
            # Save embeddings only if they exist
            if embeddings_array:
                cursor.execute("""
                    INSERT INTO content_embedding (course_material_id, ai_processed_id, chunk_text, chunk_index, chunk_type, embedding, created_at)
                    VALUES (%s, %s, %s, %s, %s, %s, %s)
                """, (material_id, ai_processed_id, extracted_text[:1000], 1, 'content', str(embeddings_array), datetime.now()))
            
            # Mark job as completed
            cursor.execute("""
                UPDATE ai_processing_job 
                SET status = 'completed', completed_at = %s 
                WHERE id = %s
            """, (datetime.now(), job_id))
            
            # Commit transaction
            conn.commit()
            logger.info(f"✅ Background processing completed for material {material_id}")
            
            # Send completion status update
            status_tracker.update_status(job_id, {
                "status": "completed",
                "stage": "completed",
                "message": "Content processing completed successfully!",
                "progress": 100,
                "material_id": material_id,
                "ai_processed_id": ai_processed_id
            })
            
        except Exception as db_error:
            # Rollback transaction on database error
            conn.rollback()
            raise DatabaseError(f"Database transaction failed: {str(db_error)}")
        finally:
            cursor.close()
            conn.close()
        
    except Exception as e:
        logger.error(f"❌ Background processing failed at stage '{processing_stage}': {str(e)}")
        
        # Send failure status update
        status_tracker.update_status(job_id, {
            "status": "failed",
            "stage": processing_stage,
            "message": f"Processing failed at {processing_stage}: {str(e)}",
            "error": str(e),
            "progress": -1
        })
        
        # Enhanced error recovery with detailed logging
        await handle_processing_failure(
            conn, job_id, material_id, ai_processed_id, 
            str(e), processing_stage
        )
        
        raise AIProcessingError(f"Background processing failed at {processing_stage}: {str(e)}")


async def retry_with_backoff(func, max_retries: int = 3, stage: str = "unknown", 
                           base_delay: float = 1.0, max_delay: float = 60.0):
    """Retry function with exponential backoff and detailed logging."""
    for attempt in range(max_retries + 1):
        try:
            result = await func()
            if attempt > 0:
                logger.info(f"✅ {stage} succeeded on attempt {attempt + 1}")
            return result
        except Exception as e:
            if attempt == max_retries:
                logger.error(f"❌ {stage} failed after {max_retries + 1} attempts: {str(e)}")
                raise e
            
            delay = min(base_delay * (2 ** attempt), max_delay)
            logger.warning(f"⚠️ {stage} failed on attempt {attempt + 1}/{max_retries + 1}, retrying in {delay}s: {str(e)}")
            await asyncio.sleep(delay)


async def handle_processing_failure(conn, job_id: str, material_id: str, 
                                  ai_processed_id: Optional[str], error_msg: str, 
                                  failed_stage: str):
    """Handle processing failure with comprehensive cleanup and recovery options."""
    try:
        # Ensure we have a database connection
        if conn is None or conn.closed:
            conn = get_db_connection()
        
        cursor = conn.cursor()
        
        # Clean up any partial data if database operations had started
        if ai_processed_id:
            try:
                # Remove partial content embeddings
                cursor.execute("""
                    DELETE FROM content_embedding WHERE ai_processed_id = %s
                """, (ai_processed_id,))
                
                # Remove partial processed content
                cursor.execute("""
                    DELETE FROM ai_processed_content WHERE id = %s
                """, (ai_processed_id,))
                
                logger.info(f"🧹 Cleaned up partial data for ai_processed_id: {ai_processed_id}")
            except Exception as cleanup_error:
                logger.error(f"❌ Failed to cleanup partial data: {str(cleanup_error)}")
        
        # Update job status with detailed error information
        cursor.execute("""
            UPDATE ai_processing_job 
            SET status = 'failed', 
                error_message = %s, 
                completed_at = %s,
                metadata = COALESCE(metadata, '{}'::jsonb) || %s::jsonb
            WHERE id = %s
        """, (
            error_msg, 
            datetime.now(), 
            json.dumps({
                "failed_stage": failed_stage,
                "failure_timestamp": datetime.now().isoformat(),
                "can_retry": failed_stage in ["text_extraction", "ai_summary", "key_concepts", "embeddings"]
            }),
            job_id
        ))
        
        conn.commit()
        cursor.close()
        
        # Log recovery suggestions
        if failed_stage in ["ai_summary", "key_concepts", "embeddings"]:
            logger.info(f"💡 Job {job_id} can be retried - AI service may be temporarily unavailable")
        elif failed_stage == "text_extraction":
            logger.info(f"💡 Job {job_id} file format issue - check file integrity")
        else:
            logger.info(f"💡 Job {job_id} requires manual intervention")
            
    except Exception as recovery_error:
        logger.error(f"❌ Error recovery failed: {str(recovery_error)}")
        # Fallback: try to at least update the job status
        try:
            if conn is None or conn.closed:
                conn = get_db_connection()
            cursor = conn.cursor()
            cursor.execute("""
                UPDATE ai_processing_job 
                SET status = 'failed', error_message = %s, completed_at = %s 
                WHERE id = %s
            """, (f"Processing failed: {error_msg}. Recovery also failed: {str(recovery_error)}", datetime.now(), job_id))
            conn.commit()
            cursor.close()
        except Exception:
            logger.critical(f"❌ Complete failure - unable to update job {job_id} status")
    finally:
        if conn and not conn.closed:
            conn.close()

def clean_extracted_text(text: str) -> str:
    """Clean extracted text to remove problematic characters for database storage."""
    if not text:
        return text
    
    # Remove null characters that cause PostgreSQL issues
    cleaned = text.replace('\x00', '')
    
    # Remove other problematic control characters but keep useful ones like newlines and tabs
    import re
    # Keep: \n (newline), \r (carriage return), \t (tab)
    # Remove: all other control characters (0x00-0x1F except \n, \r, \t)
    cleaned = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', cleaned)
    
    return cleaned

async def extract_text_from_file(file_path: str) -> str:
    """Extract text from various file formats using enhanced processors."""
    try:
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileValidationError(f"File not found: {file_path}")
        
        file_extension = file_path.suffix.lower()
        logger.info(f"🔍 Processing file type: {file_extension}")
        
        # Text files - simple read
        if file_extension == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
                logger.info(f"✅ Text file processed: {len(content)} characters")
                return clean_extracted_text(content)
        
        # PDF files - enhanced processor with metadata
        elif file_extension == '.pdf':
            if not PDF_AVAILABLE:
                raise ContentProcessingError("PDF processing not available (PyPDF2 not installed)")
            content = await enhanced_pdf_processor.process_file(str(file_path))
            return clean_extracted_text(content)
        
        # DOCX files - enhanced processor with styles and tables
        elif file_extension == '.docx':
            if not DOCX_AVAILABLE:
                raise ContentProcessingError("DOCX processing not available (python-docx not installed)")
            content = await enhanced_docx_processor.process_file(str(file_path))
            return clean_extracted_text(content)
        
        # PPTX files - enhanced processor with slide analysis
        elif file_extension == '.pptx':
            if not PPTX_AVAILABLE:
                raise ContentProcessingError("PPTX processing not available (python-pptx not installed)")
            content = await enhanced_pptx_processor.process_file(str(file_path))
            return clean_extracted_text(content)
        
        # Audio/Video files - AI transcription
        elif file_extension in ['.mp3', '.wav', '.m4a', '.flac', '.mp4', '.avi', '.mov', '.mkv', '.webm']:
            if not WHISPER_AVAILABLE:
                raise ContentProcessingError("Audio/video processing not available (whisper not installed)")
            content = await enhanced_audio_video_processor.process_file(str(file_path))
            return clean_extracted_text(content)
        
        # Other document formats that might contain text
        elif file_extension in ['.rtf', '.odt']:
            logger.warning(f"⚠️ Limited support for {file_extension} - attempting text extraction")
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                    return clean_extracted_text(content)
            except UnicodeDecodeError:
                # Try with different encodings
                for encoding in ['latin1', 'cp1252', 'iso-8859-1']:
                    try:
                        with open(file_path, 'r', encoding=encoding) as f:
                            content = f.read()
                            return clean_extracted_text(content)
                    except UnicodeDecodeError:
                        continue
                raise FileValidationError(f"Could not decode {file_extension} file with any supported encoding")
        
        # Unsupported file types
        else:
            # Last resort: try to read as text
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                    logger.warning(f"⚠️ Processed {file_extension} as plain text")
                    return clean_extracted_text(content)
            except UnicodeDecodeError:
                raise FileValidationError(
                    f"Unsupported file format: {file_extension}. "
                    f"Supported formats: PDF, DOCX, PPTX, TXT, MP3, WAV, MP4, AVI, MOV"
                )
    
    except (FileValidationError, ContentProcessingError) as e:
        # Re-raise our custom exceptions
        raise e
    except Exception as e:
        logger.error(f"❌ Text extraction failed for {file_path}: {str(e)}")
        raise ContentProcessingError(f"Failed to extract text from file: {str(e)}")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(
        app,
        host="0.0.0.0",
        port=8082,
        log_level="info"
    )